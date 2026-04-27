from __future__ import annotations
import os
import asyncio
from pathlib import Path
from typing import List, Dict, Any

import fitz  # PyMuPDF
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_openai import ChatOpenAI
from workflow_engine.workflow.registry import register
from workflow_engine.graphbuilder.graph_builder import GenericGraphBuilder
from workflow_engine.logger import get_logger
from workflow_engine.state import KBMindMapState
from workflow_engine.utils import get_project_root

log = get_logger(__name__)

MINDMAP_MARKDOWN_SYSTEM_PROMPT = """\
你是一个思维导图生成器。根据用户提供的文本内容，输出一份用于 markmap 渲染的 Markdown 思维导图。

输出要求：
- 只输出 Markdown，不要输出 JSON
- 使用 ATX 标题语法（# / ## / ###）
- 整份内容必须是一棵思维导图，不要写解释、前言、总结、代码围栏、项目符号或段落
- 标题文本要简短，突出关键概念、论点和层级关系
- 思维导图层次清晰、要点完整、用词简练
"""

QUALITY_JUDGING_CRITERIA = """\
生成思维导图请严格遵循以下四个核心原则：

1. 结构性（Structure）
使用清晰的层次结构
将复杂问题合理拆分为子节点
避免简单线性罗列，体现组织结构

2. 关联性（Connectivity）
节点之间必须有明确、合理的关系（如因果、依赖、解释等）
避免无关或错误连接
每个子节点都应与父节点存在明确语义关系

3. 落地性（Groundedness / Actionability）
所有叶子节点必须是“具体、可用的信息”，例如：事实\数据\明确结论可直接用于回答问题的内容
避免抽象、空泛或无实际信息的表述（如“进一步分析”、“深入研究”等）

4. 简洁性（Conciseness）
避免冗余、重复节点
控制分支数量，保证信息密度高
除原文表述，每个节点不可使用包含冒号和破折号的表达，一个节点只能是单一主体、概念（如“成功率”、“模型”）或短语（如“准确率90%”，“效果显著提升”）。
    
"""

SINGLE_DOC_PROMPT = """\
请根据以下文档内容生成一个总结型思维导图。
要求：
- 提取关键概念、论点和结构化要点，保留重要的数据和示例
- 最大层级深度为 {max_depth} 层
- 直接输出 Markdown 标题层级，使用 #、##、### 等markdown标记表示层级

    {QUALITY_JUDGING_CRITERIA}

=== 文档: {title} ===
{content}

直接输出最终思维导图，不要任何解释或额外文本。
""".replace("{QUALITY_JUDGING_CRITERIA}", QUALITY_JUDGING_CRITERIA)


MERGE_PROMPT = """\
我们已有多篇文档，并已分别生成了各自的思维导图。现在需要将这些导图整合为一个统一的最终思维导图。

约束
1. 先分析关系再组织结构。若文档主题或方法体系高度重叠 → 按主题融合，合并相近节点；若主题差异明显 → 保持为并列主分支，不强行统一抽象
2. 结构要求
   主分支使用“主题”命名，而不是文档名
   相似内容需合并去重
   每篇文档的独有信息必须保留，避免丢失或过度概括
3. 保留重要的数据和示例
4. 直接输出 Markdown 标题层级，使用 #、##、### 等markdown标记表示层级
5. 最大层级深度为：{max_depth}

    {QUALITY_JUDGING_CRITERIA}

不输出任何解释，仅输出最终思维导图

已有的文档思维导图如下：
{document_mindmaps}
""".replace("{QUALITY_JUDGING_CRITERIA}", QUALITY_JUDGING_CRITERIA)


MINDMAP_SECTION = """\
=== 文档思维导图: {title} ===
{mindmap_markdown}
"""

# Try importing office libraries
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

@register("kb_mindmap")
def create_kb_mindmap_graph() -> GenericGraphBuilder:
    """
    Workflow for Knowledge Base MindMap Generation
    Steps:
    1. Parse uploaded files (PDF/Office)
    2. Analyze content structure using LLM
    3. Generate Mermaid mindmap syntax using LLM
    """
    builder = GenericGraphBuilder(state_model=KBMindMapState, entry_point="_start_")

    def _strip_markdown_code_fence(text: str) -> str:
        cleaned = str(text or "").strip()
        if "```" not in cleaned:
            return cleaned
        lines = cleaned.splitlines()
        in_code_block = False
        code_lines: List[str] = []
        for line in lines:
            if line.strip().startswith("```"):
                in_code_block = not in_code_block
                continue
            if in_code_block:
                code_lines.append(line)
        return "\n".join(code_lines).strip() or cleaned

    def _normalize_markdown_headings(text: str, max_depth: int) -> str:
        lines: List[str] = []
        for raw_line in _strip_markdown_code_fence(text).splitlines():
            line = raw_line.rstrip()
            stripped = line.lstrip()
            if not stripped:
                continue
            if stripped.startswith("#"):
                hashes = len(stripped) - len(stripped.lstrip("#"))
                title = stripped[hashes:].strip()
                if title:
                    level = max(1, min(hashes, max_depth))
                    lines.append(f"{'#' * level} {title}")
            elif lines:
                lines.append(f"{'#' * min(max_depth, 3)} {stripped.strip('-* ')}")
        return "\n".join(lines).strip()

    async def _run_docmerge_markdown_prompt(
        state: KBMindMapState,
        prompt: str,
        *,
        temperature: float,
    ) -> str:
        llm = ChatOpenAI(
            openai_api_base=state.request.chat_api_url,
            openai_api_key=state.request.api_key,
            model_name=state.request.model,
            temperature=temperature,
        )
        response = await llm.ainvoke([
            SystemMessage(content=MINDMAP_MARKDOWN_SYSTEM_PROMPT),
            HumanMessage(content=prompt),
        ])
        return str(response.content or "")

    def _single_doc_prompt(filename: str, content: str, max_depth: int, language: str) -> str:
        return SINGLE_DOC_PROMPT.format(title=filename, content=content, max_depth=max_depth)

    def _merge_prompt(document_trees: List[Dict[str, str]], max_depth: int, language: str) -> str:
        tree_sections = "\n\n".join(
            MINDMAP_SECTION.format(title=item["filename"], mindmap_markdown=item["markdown"])
            for item in document_trees
        )
        return MERGE_PROMPT.format(document_mindmaps=tree_sections, max_depth=max_depth)

    def _start_(state: KBMindMapState) -> KBMindMapState:
        # Ensure request fields
        if not state.request.file_ids:
            state.request.file_ids = []

        # Initialize output directory
        if not state.result_path:
            project_root = get_project_root()
            import time
            ts = int(time.time())
            email = getattr(state.request, 'email', 'default')
            output_dir = project_root / "outputs" / "kb_outputs" / email / f"{ts}_mindmap"
            output_dir.mkdir(parents=True, exist_ok=True)
            state.result_path = str(output_dir)

        state.file_contents = []
        state.content_structure = ""
        state.mermaid_code = ""
        state.mindmap_svg_path = ""
        return state

    async def parse_files_node(state: KBMindMapState) -> KBMindMapState:
        """
        Parse all files and extract content
        """
        files = state.request.file_ids
        if not files:
            state.file_contents = []
            return state

        async def process_file(file_path: str) -> Dict[str, Any]:
            file_path_obj = Path(file_path)
            filename = file_path_obj.name

            if not file_path_obj.exists():
                return {
                    "filename": filename,
                    "content": f"[Error: File not found {file_path}]"
                }

            suffix = file_path_obj.suffix.lower()
            raw_content = ""

            try:
                # PDF
                if suffix == ".pdf":
                    try:
                        doc = fitz.open(file_path)
                        text = ""
                        for page in doc:
                            text += page.get_text() + "\n"
                        raw_content = text
                    except Exception as e:
                        raw_content = f"[Error parsing PDF: {e}]"

                # Word
                elif suffix in [".docx", ".doc"]:
                    if Document is None:
                         raw_content = "[Error: python-docx not installed]"
                    else:
                        try:
                            doc = Document(file_path)
                            raw_content = "\n".join([p.text for p in doc.paragraphs])
                        except Exception as e:
                             raw_content = f"[Error parsing Docx: {e}]"

                # PPT
                elif suffix in [".pptx", ".ppt"]:
                    if Presentation is None:
                        raw_content = "[Error: python-pptx not installed]"
                    else:
                        try:
                            prs = Presentation(file_path)
                            text = ""
                            for i, slide in enumerate(prs.slides):
                                text += f"--- Slide {i+1} ---\n"
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text += shape.text + "\n"
                            raw_content = text
                        except Exception as e:
                            raw_content = f"[Error parsing PPT: {e}]"

                else:
                    try:
                        with open(file_path, "r", encoding="utf-8") as f:
                            raw_content = f.read()
                    except:
                        raw_content = "[Unsupported file type]"

            except Exception as e:
                 raw_content = f"[Parse Error: {e}]"

            # Truncate content
            truncated_content = raw_content[:50000] if len(raw_content) > 50000 else raw_content

            return {
                "filename": filename,
                "content": truncated_content
            }

        # Run in parallel
        tasks = [process_file(f) for f in files]
        results = await asyncio.gather(*tasks)

        state.file_contents = results
        return state

    async def analyze_structure_node(state: KBMindMapState) -> KBMindMapState:
        """Skip structure analysis: single-pass generation happens downstream."""
        state.content_structure = "merged"
        return state

    async def generate_mermaid_node(state: KBMindMapState) -> KBMindMapState:
        """Generate a mindmap with the docmerge strategy: per-document trees, then merge."""
        if not state.file_contents:
            state.mermaid_code = "# Error\n## No content available"
            return state

        language = state.request.language
        max_depth = state.request.max_depth

        try:
            per_doc_maps: List[Dict[str, str]] = []
            for item in state.file_contents:
                filename = str(item.get("filename") or "未命名文档")
                content = str(item.get("content") or "").strip()
                if not content:
                    continue
                markdown = await _run_docmerge_markdown_prompt(
                    state,
                    _single_doc_prompt(filename, content, max_depth, language),
                    temperature=0.3,
                )
                normalized = _normalize_markdown_headings(markdown, max_depth=max_depth)
                if not normalized:
                    continue
                per_doc_maps.append({
                    "filename": filename,
                    "markdown": normalized,
                })

            if not per_doc_maps:
                state.mermaid_code = "# Error\n## No content available"
            elif len(per_doc_maps) == 1:
                state.mermaid_code = per_doc_maps[0]["markdown"]
            else:
                merged_markdown = await _run_docmerge_markdown_prompt(
                    state,
                    _merge_prompt(per_doc_maps, max_depth, language),
                    temperature=0.3,
                )
                state.mermaid_code = _normalize_markdown_headings(merged_markdown, max_depth=max_depth)
                if not state.mermaid_code:
                    state.mermaid_code = "# Error\n## Generation failed"
            if per_doc_maps:
                state.content_structure = "\n\n".join(
                    f"=== 文档思维导图: {item['filename']} ===\n{item['markdown']}"
                    for item in per_doc_maps
                )
        except Exception as e:
            log.error(f"DocMerge mindmap generation failed: {e}")
            state.mermaid_code = f"# Error\n## {str(e)}"

        # Save mermaid code to file
        try:
            mermaid_path = Path(state.result_path) / "mindmap.mmd"
            mermaid_path.parent.mkdir(parents=True, exist_ok=True)
            mermaid_path.write_text(state.mermaid_code, encoding="utf-8")
            log.info(f"Mermaid code saved to: {mermaid_path}")
        except Exception as e:
            log.error(f"Failed to save mermaid code: {e}")

        return state

    nodes = {
        "_start_": _start_,
        "parse_files": parse_files_node,
        "analyze_structure": analyze_structure_node,
        "generate_mermaid": generate_mermaid_node,
        "_end_": lambda s: s
    }

    edges = [
        ("_start_", "parse_files"),
        ("parse_files", "analyze_structure"),
        ("analyze_structure", "generate_mermaid"),
        ("generate_mermaid", "_end_")
    ]

    builder.add_nodes(nodes).add_edges(edges)
    return builder
