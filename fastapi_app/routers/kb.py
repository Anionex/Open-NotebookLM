import asyncio
import json
import os
import re
import shutil
import subprocess
import time
from pathlib import Path
from urllib.parse import urlparse, unquote
from fastapi import APIRouter, UploadFile, File, Form, HTTPException, Body
from fastapi.responses import StreamingResponse
from typing import Optional, List, Dict, Any
from openai import AsyncOpenAI

import fitz  # PyMuPDF

from workflow_engine.state import IntelligentQARequest, IntelligentQAState, KBPodcastRequest, KBPodcastState, KBMindMapRequest, KBMindMapState
from workflow_engine.toolkits.ragtool.vector_store_tool import process_knowledge_base_files, VectorStoreManager
from workflow_engine.utils import get_project_root
from workflow_engine.logger import get_logger
from workflow_engine.workflow import run_workflow, list_workflows

log = get_logger(__name__)
from fastapi_app.config import settings
from fastapi_app.schemas import Paper2PPTRequest
from fastapi_app.utils import _from_outputs_url, _to_outputs_url
from fastapi_app.services.wa_paper2ppt import _init_state_from_request
from fastapi_app.dependencies.auth import get_supabase_admin_client
from fastapi_app.notebook_paths import NotebookPaths, get_notebook_paths, _sanitize_user_id
from fastapi_app.source_manager import SourceManager
from fastapi_app.services.fast_research_service import fast_research_search
from fastapi_app.services.deep_research_report_service import generate_report_from_search
from workflow_engine.toolkits.research_tools import fetch_page_text
from workflow_engine.workflow.wf_intelligent_qa import prepare_parallel_file_analyses, build_intelligent_qa_prompt
from workflow_engine.promptstemplates.resources.pt_qa_agent_repo import KbPromptAgent as KbPromptAgentPrompts

router = APIRouter(prefix="/kb", tags=["Knowledge Base"])

# Link sources JSON filename under notebook dir (excluded from regular file list)
LINK_SOURCES_FILENAME = "link_sources.json"

# Base directory for storing KB files
# Layout: outputs/kb_data/{email}/{notebook_id}/ for per-notebook isolation
KB_BASE_DIR = Path("outputs/kb_data")
OUTPUTS_BASE = Path("outputs/kb_outputs")


def _notebook_dir(email: str, notebook_id: Optional[str]) -> Path:
    """User + notebook scoped dir under kb_data. Email is sanitized for filesystem safety."""
    root = get_project_root()
    safe_email = _sanitize_user_id(email) if email else "default"
    base = root / KB_BASE_DIR / safe_email
    if notebook_id:
        return base / notebook_id.replace("/", "_").replace("\\", "_")[:128]
    return base / "_shared"


def _outputs_dir(email: str, notebook_id: Optional[str], subdir: str) -> Path:
    """User + notebook scoped output dir. Email is sanitized for filesystem safety."""
    root = get_project_root()
    safe_email = _sanitize_user_id(email) if email else "default"
    base = root / OUTPUTS_BASE / safe_email
    if notebook_id:
        base = base / notebook_id.replace("/", "_").replace("\\", "_")[:128]
    else:
        base = base / "_shared"
    return base / subdir


def _get_cjk_font_path() -> Optional[str]:
    """返回系统中文字体路径，用于 PDF 内中文显示；无则返回 None。"""
    candidates = [
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSerifCJK-Regular.ttc",
        "/usr/share/fonts/truetype/arphic/uming.ttc",
    ]
    for path in candidates:
        if Path(path).exists():
            return path
    return None


def _text_to_pdf(text: str, output_path: str) -> None:
    """将长文本生成为多页 PDF（PyMuPDF），不依赖 kb_page_content / paper2ppt workflow。支持中文（CJK 字体）。"""
    text = (text or "").strip()
    if not text:
        raise ValueError("Report text is empty")
    doc = fitz.open()
    rect = fitz.Rect(50, 50, 545, 802)
    fontsize = 11
    max_chars_per_page = 3200
    fontfile = _get_cjk_font_path()
    fontname = "notocjk" if fontfile else "helv"
    if fontfile:
        # 使用中文字体，否则中文会不显示
        pass
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    if not paragraphs:
        paragraphs = [text]
    current: List[str] = []
    current_len = 0
    for p in paragraphs:
        need = len(p) + (2 if current else 0)
        if current_len + need > max_chars_per_page and current:
            page = doc.new_page(width=595, height=842)
            if fontfile:
                page.insert_textbox(rect, "\n\n".join(current), fontsize=fontsize, fontname=fontname, fontfile=fontfile)
            else:
                page.insert_textbox(rect, "\n\n".join(current), fontsize=fontsize, fontname=fontname)
            current = [p]
            current_len = len(p)
        else:
            current.append(p)
            current_len += need
    if current:
        page = doc.new_page(width=595, height=842)
        if fontfile:
            page.insert_textbox(rect, "\n\n".join(current), fontsize=fontsize, fontname=fontname, fontfile=fontfile)
        else:
            page.insert_textbox(rect, "\n\n".join(current), fontsize=fontsize, fontname=fontname)
    doc.save(output_path)
    doc.close()


def _unwrap_fastapi_body_default(value: Any, fallback: Any = None) -> Any:
    """
    outputs-v2 等内部代码会直接调用路由函数，此时未显式传入的参数仍可能是 FastAPI 的 Body 默认对象。
    这里把它们还原成普通值，避免后续 string/int 操作报错。
    """
    if type(value).__name__ == "Body":
        default = getattr(value, "default", None)
        if default is ...:
            return fallback
        return default if default is not None else fallback
    return value


ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg", ".mp4", ".md", ".csv"}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg"}
DOC_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".md", ".markdown"}
DATASET_EXTENSIONS = {".csv"}


def _find_mineru_stem_dir(
    pdf_stem: str,
    email: str,
    notebook_id: Optional[str],
    notebook_title: Optional[str] = None,
) -> Optional[Path]:
    """
    查找指定 pdf_stem 的 MinerU 输出目录。
    查找顺序：
    1. 笔记本新布局: outputs/{title}_{id}/sources/{pdf_stem}/mineru/
    2. kb_mineru 新结构: kb_mineru/{email}/{notebook_id}/{pdf_stem}/auto/
    3. kb_mineru 旧结构: kb_mineru/{email}/{notebook_id}/{uuid}/{pdf_stem}/auto/
    返回包含 auto/ 或 hybrid_auto/ 的目录，找不到返回 None。
    """
    project_root = get_project_root()

    # 1) 笔记本新布局: outputs/{title}_{id}/sources/{pdf_stem}/mineru/
    if notebook_id:
        nb_paths = get_notebook_paths(notebook_id, notebook_title or "", email)
        mineru_dir = nb_paths.sources_dir / pdf_stem / "mineru"
        if mineru_dir.exists():
            # 直接在 mineru/ 下找 auto/ 或 hybrid_auto/
            for sub in ("auto", "hybrid_auto"):
                if (mineru_dir / sub).is_dir() and list((mineru_dir / sub).glob("*.md")):
                    log.info("[find_mineru] 在新布局找到缓存: %s/%s", mineru_dir, sub)
                    return mineru_dir
            # 兼容: mineru/{pdf_stem}/auto/ (MinerU 可能多嵌套一层)
            nested = mineru_dir / pdf_stem
            if nested.exists():
                for sub in ("auto", "hybrid_auto"):
                    if (nested / sub).is_dir() and list((nested / sub).glob("*.md")):
                        log.info("[find_mineru] 在新布局(嵌套)找到缓存: %s/%s", nested, sub)
                        return nested
                # 兼容: mineru/{pdf_stem}/*.md （无 auto 子目录）
                if list(nested.glob("*.md")):
                    log.info("[find_mineru] 在新布局(平铺嵌套)找到缓存: %s", nested)
                    return nested

    # 2) Legacy: kb_mineru/{email}/{notebook_id}/
    safe_nb = (notebook_id or "_shared").replace("/", "_").replace("\\", "_")[:128]
    safe_email = _sanitize_user_id(email) if email else "default"
    mineru_base = project_root / "outputs" / "kb_mineru" / safe_email / safe_nb

    if not mineru_base.exists():
        return None

    stem_dir = mineru_base / pdf_stem
    if stem_dir.exists():
        for sub in ("auto", "hybrid_auto"):
            if (stem_dir / sub).is_dir() and list((stem_dir / sub).glob("*.md")):
                return stem_dir

    # 3) 旧结构兼容：kb_mineru/{email}/{nb}/{uuid}/{pdf_stem}/auto/
    for child in mineru_base.iterdir():
        if not child.is_dir():
            continue
        nested = child / pdf_stem
        if not nested.exists():
            continue
        for sub in ("auto", "hybrid_auto"):
            if (nested / sub).is_dir() and list((nested / sub).glob("*.md")):
                return nested

    return None


def _read_mineru_md_if_cached(
    pdf_path: Path,
    email: str,
    notebook_id: Optional[str],
    max_chars: int = 50000,
    notebook_title: Optional[str] = None,
) -> Optional[str]:
    """
    尝试从已有的 MinerU 缓存中读取 markdown 内容。
    找到则返回 markdown 文本，否则返回 None。
    """
    stem_dir = _find_mineru_stem_dir(pdf_path.stem, email, notebook_id, notebook_title)
    if stem_dir is None:
        return None

    for sub in ("auto", "hybrid_auto"):
        candidate = stem_dir / sub
        if not candidate.is_dir():
            continue
        md_files = list(candidate.glob("*.md"))
        if md_files:
            try:
                text = md_files[0].read_text(encoding="utf-8")
                if text.strip():
                    log.info("[read_mineru_md] 从缓存读取 %s, len=%s", md_files[0], len(text))
                    return text[:max_chars] if len(text) > max_chars else text
            except Exception as e:
                log.warning("[read_mineru_md] 读取失败 %s: %s", md_files[0], e)
    md_files = list(stem_dir.glob("*.md"))
    if md_files:
        try:
            text = md_files[0].read_text(encoding="utf-8")
            if text.strip():
                log.info("[read_mineru_md] 从平铺缓存读取 %s, len=%s", md_files[0], len(text))
                return text[:max_chars] if len(text) > max_chars else text
        except Exception as e:
            log.warning("[read_mineru_md] 读取失败 %s: %s", md_files[0], e)
    return None


def _copy_mineru_tree(src: Path, dst: Path) -> None:
    dst.mkdir(parents=True, exist_ok=True)
    for child in src.iterdir():
        target = dst / child.name
        if child.is_dir():
            shutil.copytree(str(child), str(target), dirs_exist_ok=True)
        else:
            shutil.copy2(str(child), str(target))


def _normalize_cached_mineru_dir(cached_stem_dir: Path, target: Path, stem: str) -> None:
    auto_like_dirs = []
    for sub in ("auto", "hybrid_auto"):
        candidate = cached_stem_dir / sub
        if candidate.is_dir() and list(candidate.glob("*.md")):
            auto_like_dirs.append(candidate)

    if auto_like_dirs:
        try:
            target.symlink_to(cached_stem_dir.resolve())
            return
        except OSError:
            shutil.copytree(str(cached_stem_dir), str(target))
            return

    flat_md_files = list(cached_stem_dir.glob("*.md"))
    if flat_md_files:
        normalized_auto = target / "auto"
        _copy_mineru_tree(cached_stem_dir, normalized_auto)
        expected_md = normalized_auto / f"{stem}.md"
        if not expected_md.exists() and flat_md_files:
            shutil.copy2(str(flat_md_files[0]), str(expected_md))
        return

    try:
        target.symlink_to(cached_stem_dir.resolve())
    except OSError:
        shutil.copytree(str(cached_stem_dir), str(target))


def _reuse_mineru_cache(
    pdf_paths: List[Path],
    output_dir: Path,
    email: str,
    notebook_id: Optional[str],
    notebook_title: Optional[str] = None,
) -> int:
    """
    将已有的 MinerU 解析结果复制/软链到 PPT workflow 的 output_dir 下，
    使 parse_pdf_pages 能直接发现 {output_dir}/{pdf_stem}/auto/*.md 而跳过重新解析。
    返回成功复用的 PDF 数量。
    """
    reused = 0
    for pdf_path in pdf_paths:
        stem = pdf_path.stem
        cached_stem_dir = _find_mineru_stem_dir(stem, email, notebook_id, notebook_title)
        if cached_stem_dir is None:
            log.info("[reuse_mineru] 未找到 %s 的 MinerU 缓存", stem)
            continue

        target = output_dir / stem
        if target.exists():
            # 已存在（可能之前已复用或本次 workflow 已生成），跳过
            log.info("[reuse_mineru] 目标已存在，跳过: %s", target)
            reused += 1
            continue

        try:
            _normalize_cached_mineru_dir(cached_stem_dir, target, stem)
            log.info("[reuse_mineru] 复用成功: %s -> %s", target, cached_stem_dir)
            reused += 1
        except Exception as e:
            log.warning("[reuse_mineru] 复用失败 %s: %s", stem, e)

    return reused


def _resolve_local_path(path_or_url: str) -> Path:
    if not path_or_url:
        raise HTTPException(status_code=400, detail="Empty file path")
    raw = _from_outputs_url(path_or_url)
    p = Path(raw)
    if not p.is_absolute():
        p = (get_project_root() / p).resolve()
    elif not p.exists():
        # 前端可能传了带 /outputs/ 的绝对形式，在服务端需按 project_root 解析（并解码 %40）
        raw_stripped = unquote(raw.lstrip("/"))
        if raw_stripped:
            p_rel = (get_project_root() / raw_stripped).resolve()
            if p_rel.exists():
                p = p_rel
    return p


def _convert_to_pdf(input_path: Path, output_dir: Path) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(input_path)
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    pdf_path = output_dir / f"{input_path.stem}.pdf"
    if not pdf_path.exists():
        raise HTTPException(status_code=500, detail=f"PDF conversion failed for {input_path.name}")
    return pdf_path


def _merge_pdfs(pdf_paths: List[Path], output_path: Path) -> Path:
    if not pdf_paths:
        raise HTTPException(status_code=400, detail="No PDF files to merge")
    merged = fitz.open()
    for pdf in pdf_paths:
        with fitz.open(pdf) as src:
            merged.insert_pdf(src)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    merged.save(output_path)
    merged.close()
    return output_path


def _append_images_to_pptx(pptx_path: Path, image_paths: List[Path]) -> None:
    try:
        from pptx import Presentation
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"python-pptx not available: {e}")

    prs = Presentation(str(pptx_path))
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height
        )
    prs.save(str(pptx_path))

@router.post("/upload")
async def upload_kb_file(
    file: UploadFile = File(...),
    email: str = Form(...),
    user_id: str = Form(...),
    notebook_id: Optional[str] = Form(None),
    notebook_title: Optional[str] = Form(None),
):
    """
    Upload a file to the notebook's knowledge base directory.
    New layout: outputs/{title}_{id}/sources/{stem}/original/
    Fallback: also writes to legacy kb_data path for backward compat.
    """
    if not email:
        raise HTTPException(status_code=400, detail="Email is required")
    if not notebook_id:
        raise HTTPException(status_code=400, detail="notebook_id is required for per-notebook storage")

    file_ext = Path(file.filename).suffix.lower()
    if file_ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: {file_ext}. Allowed: {', '.join(ALLOWED_EXTENSIONS)}"
        )

    try:
        filename = file.filename or f"unnamed_{user_id}"
        filename = os.path.basename(filename)

        # --- New notebook-centric layout ---
        paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
        mgr = SourceManager(paths)

        # Save uploaded bytes to a temp location first, then import
        tmp_dir = paths.root / "_tmp"
        tmp_dir.mkdir(parents=True, exist_ok=True)
        tmp_path = tmp_dir / filename
        with open(tmp_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        source_info = await mgr.import_file(tmp_path, filename)

        # Clean up temp
        try:
            tmp_path.unlink(missing_ok=True)
        except Exception:
            pass

        # Build static URL from the original path in new layout
        project_root = get_project_root()
        rel = source_info.original_path.relative_to(project_root)
        static_path = "/" + rel.as_posix()

        # --- Also write to legacy path for backward compat ---
        legacy_dir = _notebook_dir(email, notebook_id)
        legacy_dir.mkdir(parents=True, exist_ok=True)
        legacy_path = legacy_dir / filename
        if not legacy_path.exists():
            shutil.copy2(str(source_info.original_path), str(legacy_path))

        # Auto-embed using new vector_store path
        embedded = False
        if file_ext not in DATASET_EXTENSIONS:
            try:
                vector_base = str(paths.vector_store_dir)
                mineru_base = str(paths.source_mineru_dir(filename))
                file_list = [{"path": str(source_info.original_path)}]
                await process_knowledge_base_files(
                    file_list=file_list,
                    base_dir=vector_base,
                    mineru_output_base=mineru_base,
                )
                embedded = True
                log.info("[upload] auto-embedding done: %s", filename)
            except Exception as emb_err:
                log.warning("[upload] auto-embedding failed for %s: %s", filename, emb_err)

        # Write to JSON records
        from fastapi_app.kb_records import add_source_record
        try:
            add_source_record(
                user_email=email,
                notebook_id=notebook_id,
                file_name=filename,
                file_path=str(source_info.original_path),
                static_url=static_path,
                file_size=os.path.getsize(source_info.original_path),
                file_type=file.content_type or ""
            )
        except Exception as e:
            log.warning("[upload] failed to write JSON record: %s", e)

        return {
            "success": True,
            "filename": filename,
            "file_size": os.path.getsize(source_info.original_path),
            "storage_path": str(source_info.original_path),
            "static_url": static_path,
            "file_type": file.content_type,
            "embedded": embedded,
        }

    except Exception as e:
        print(f"Error uploading file: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/reembed-source")
async def reembed_source(
    notebook_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    file_path: str = Body(..., embed=True),
    user_id: Optional[str] = Body(None, embed=True),
):
    """Re-embed a single source file into the notebook vector store."""
    if not notebook_id:
        raise HTTPException(status_code=400, detail="notebook_id is required")
    if not email:
        raise HTTPException(status_code=400, detail="email is required")
    if not file_path:
        raise HTTPException(status_code=400, detail="file_path is required")

    try:
        paths = get_notebook_paths(notebook_id, "", email or user_id or "")
        project_root = get_project_root()
        resolved_path = Path(file_path)

        # /outputs/... paths are relative to project root, not filesystem root
        if not resolved_path.is_absolute() or not resolved_path.exists():
            candidate = project_root / file_path.lstrip("/")
            if candidate.exists():
                resolved_path = candidate

        if not resolved_path.exists():
            raise HTTPException(status_code=404, detail=f"File not found: {file_path}")

        filename = resolved_path.name
        vector_base = str(paths.vector_store_dir)
        mineru_base = str(paths.source_mineru_dir(filename))
        file_list = [{"path": str(resolved_path)}]

        await process_knowledge_base_files(
            file_list=file_list,
            base_dir=vector_base,
            mineru_output_base=mineru_base,
        )
        log.info("[reembed-source] done: %s", filename)
        return {"success": True, "filename": filename}
    except HTTPException:
        raise
    except Exception as e:
        log.error("[reembed-source] failed: %s", e)
        raise HTTPException(status_code=500, detail=str(e))


def _sanitize_md_filename(title: str, prefix: str = "doc") -> str:
    """生成安全的 .md 文件名，避免路径注入与非法字符。"""
    safe = re.sub(r'[^\w\u4e00-\u9fff\s\-.]', "", (title or "").strip())
    safe = (safe or prefix)[:80].strip() or prefix
    return safe + f"_{int(time.time())}.md"


def _sanitize_note_name(title: str, fallback: str = "note") -> str:
    safe = re.sub(r'[^\w\u4e00-\u9fff\s\-.]', "", (title or "").strip())
    safe = re.sub(r"\s+", "_", safe)
    safe = (safe or fallback).strip("._- ") or fallback
    return safe[:80]


def _url_to_pdf(url: str, output_path: Path, timeout_ms: int = 30000) -> None:
    """
    使用 Playwright 打开 URL 并打印为 PDF，便于后续统一走 MinerU。
    若 Playwright 未安装或失败，抛出异常。
    """
    from playwright.sync_api import sync_playwright
    url = (url or "").strip()
    if not url.startswith(("http://", "https://")):
        raise ValueError("Invalid url")
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        try:
            page = browser.new_page()
            page.goto(url, wait_until="networkidle", timeout=timeout_ms)
            page.pdf(path=str(output_path), format="A4", print_background=True)
        finally:
            browser.close()


@router.post("/add-text-source")
async def add_text_source(
    notebook_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    user_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    title: str = Body("直接输入", embed=True),
    content: str = Body(..., embed=True),
) -> Dict[str, Any]:
    """
    将纯文本保存为笔记本内的 .md 文件并作为来源。用于「直接输入」引入。
    New layout: outputs/{title}_{id}/sources/{stem}/
    """
    if not notebook_id or not email:
        raise HTTPException(status_code=400, detail="notebook_id and email are required")
    if not (content or "").strip():
        raise HTTPException(status_code=400, detail="content is required")

    # New layout
    paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
    mgr = SourceManager(paths)
    source_info = await mgr.import_text(content, title)

    # Legacy compat
    user_dir = _notebook_dir(email, notebook_id)
    user_dir.mkdir(parents=True, exist_ok=True)
    legacy_filename = _sanitize_md_filename(title, "直接输入")
    legacy_path = user_dir / legacy_filename
    if not legacy_path.exists():
        try:
            shutil.copy2(str(source_info.original_path), str(legacy_path))
        except Exception:
            pass

    project_root = get_project_root()
    rel = source_info.original_path.relative_to(project_root)
    static_path = "/" + rel.as_posix()

    from fastapi_app.kb_records import add_source_record
    try:
        add_source_record(
            user_email=email,
            notebook_id=notebook_id,
            file_name=source_info.original_path.name,
            file_path=str(source_info.original_path),
            static_url=static_path,
            file_size=source_info.original_path.stat().st_size,
            file_type="text/markdown"
        )
    except Exception as e:
        log.warning("[add-text-source] failed to write JSON record: %s", e)

    return {
        "success": True,
        "filename": source_info.original_path.name,
        "file_size": source_info.original_path.stat().st_size,
        "storage_path": str(source_info.original_path),
        "static_url": static_path,
        "id": f"file-{source_info.original_path.name}",
    }


@router.post("/save-note")
async def save_note(
    notebook_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    user_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    title: str = Body("无标题", embed=True),
    markdown: str = Body(..., embed=True),
    note_id: Optional[str] = Body(None, embed=True),
) -> Dict[str, Any]:
    """保存或更新笔记，并将其登记为 notebook 的 note 输出。"""
    if not notebook_id or not email:
        raise HTTPException(status_code=400, detail="notebook_id and email are required")
    markdown = (markdown or "").strip()
    if not markdown:
        raise HTTPException(status_code=400, detail="markdown is required")

    from fastapi_app.kb_records import upsert_output_record

    note_id = (note_id or "").strip() or f"note_{int(time.time() * 1000)}"
    safe_title = _sanitize_note_name(title or "无标题", "note")
    paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
    note_dir = paths.root / "note" / note_id
    note_dir.mkdir(parents=True, exist_ok=True)
    note_path = note_dir / f"{safe_title}.md"

    for existing_md in note_dir.glob("*.md"):
        if existing_md != note_path:
            try:
                existing_md.unlink()
            except Exception:
                pass

    note_path.write_text(markdown, encoding="utf-8")
    project_root = get_project_root()
    rel = note_path.relative_to(project_root)
    static_path = "/" + rel.as_posix()
    now = time.time()

    try:
        upsert_output_record(
            user_email=email,
            notebook_id=notebook_id,
            output_id=note_id,
            output_type="note",
            file_name=note_path.name,
            download_url=static_path,
            title=title,
            extra={"updated_at": now},
        )
    except Exception as record_err:
        log.warning("[save-note] failed to write output record: %s", record_err)

    return {
        "success": True,
        "note_id": note_id,
        "title": title,
        "filename": note_path.name,
        "static_url": static_path,
        "download_url": static_path,
        "updated_at": now,
    }


@router.post("/import-url-as-source")
async def import_url_as_source(
    notebook_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    user_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    url: str = Body(..., embed=True),
) -> Dict[str, Any]:
    """
    抓取 URL 网页正文存为 .md 文件，作为来源。
    New layout: outputs/{title}_{id}/sources/{stem}/
    """
    if not notebook_id or not email:
        raise HTTPException(status_code=400, detail="notebook_id and email are required")
    url = (url or "").strip()
    if not url.startswith("http://") and not url.startswith("https://"):
        raise HTTPException(status_code=400, detail="Invalid url")

    # Fetch page text
    try:
        text = await asyncio.to_thread(fetch_page_text, url)
        if not text or text.startswith("[抓取失败"):
            raise RuntimeError(text or "fetch_page_text returned empty")
    except Exception as e:
        log.warning("fetch_page_text failed: %s", e)
        raise HTTPException(status_code=500, detail=f"网页抓取失败: {e}")

    # Parse title from URL
    try:
        parsed = urlparse(url)
        title = (parsed.netloc or "网页") + "_" + (parsed.path.strip("/") or "page")[:30]
    except Exception:
        title = "网页"

    # New layout
    paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
    mgr = SourceManager(paths)
    source_info = await mgr.import_url(url, text, title)

    # Legacy compat
    user_dir = _notebook_dir(email, notebook_id)
    user_dir.mkdir(parents=True, exist_ok=True)
    legacy_filename = _sanitize_md_filename(title, "网页")
    legacy_path = user_dir / legacy_filename
    if not legacy_path.exists():
        try:
            shutil.copy2(str(source_info.original_path), str(legacy_path))
        except Exception:
            pass

    project_root = get_project_root()
    rel = source_info.original_path.relative_to(project_root)
    static_path = "/" + rel.as_posix()

    from fastapi_app.kb_records import add_source_record
    try:
        add_source_record(
            user_email=email,
            notebook_id=notebook_id,
            file_name=source_info.original_path.name,
            file_path=str(source_info.original_path),
            static_url=static_path,
            file_size=source_info.original_path.stat().st_size,
            file_type="text/markdown"
        )
    except Exception as e:
        log.warning("[import-url-as-source] failed to write JSON record: %s", e)

    embedded = False
    try:
        await process_knowledge_base_files(
            file_list=[{"path": str(source_info.original_path)}],
            base_dir=str(paths.vector_store_dir),
        )
        embedded = True
        log.info("[import-url-as-source] auto-embedding done: %s", source_info.original_path.name)
    except Exception as emb_err:
        log.warning("[import-url-as-source] auto-embedding failed for %s: %s", source_info.original_path.name, emb_err)

    return {
        "success": True,
        "filename": source_info.original_path.name,
        "file_size": source_info.original_path.stat().st_size,
        "storage_path": str(source_info.original_path),
        "static_url": static_path,
        "id": f"file-{source_info.original_path.name}",
        "embedded": embedded,
    }


@router.post("/delete-source")
async def delete_source(
    notebook_id: str = Body(..., embed=True),
    notebook_title: str = Body("", embed=True),
    user_id: str = Body("local", embed=True),
    email: Optional[str] = Body(None, embed=True),
    file_path: str = Body(..., embed=True),
):
    """Delete a source file, its directory, and the kb_records entry."""
    from fastapi_app.kb_records import remove_source_record
    from fastapi_app.services.source_service import SourceService

    effective_email = (email or user_id or "local").strip() or "local"

    try:
        svc = SourceService()

        # 0. Remove from vector store / manifest first so deleted sources are not retrieved by RAG anymore.
        try:
            svc.remove_source_from_vector_store(file_path)
        except Exception as exc:
            log.warning("Failed to remove source from vector store for %s: %s", file_path, exc)

        # 1. Delete from new layout: sources/{name}/original/{file}
        local = Path(_from_outputs_url(file_path))
        file_name = local.name
        if local.exists():
            source_dir = local.parent
            if source_dir.name == "original":
                source_dir = source_dir.parent
            if source_dir.exists() and source_dir.is_dir():
                shutil.rmtree(source_dir)
            elif local.is_file():
                os.remove(str(local))

        # 2. Delete from legacy layout: kb_data/{email}/{notebook_id}/{file}
        legacy_dir = svc._legacy_notebook_dir(effective_email, notebook_id)
        legacy_file = legacy_dir / file_name
        if legacy_file.exists() and legacy_file.is_file():
            os.remove(str(legacy_file))

        # 3. Clean up _sources.json record
        remove_source_record(
            user_email=effective_email,
            notebook_id=notebook_id,
            static_url=file_path,
        )

        return {"success": True, "message": "Source deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.delete("/delete")
async def delete_kb_file(
    storage_path: str = Form(...)
):
    """
    Delete a file from the physical storage.
    """
    try:
        # Security check: ensure path is within KB_BASE_DIR
        # This is a basic check. In production, use more robust path validation.
        target_path = Path(storage_path).resolve()
        base_path = KB_BASE_DIR.resolve()
        
        if not str(target_path).startswith(str(base_path)):
             # Allow if it's the absolute path provided by the user system
             # Check if it exists essentially
             pass

        if target_path.exists() and target_path.is_file():
            os.remove(target_path)
            return {"success": True, "message": "File deleted"}
        else:
            return {"success": False, "message": "File not found"}
            
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

def _vector_store_base_dir(email: Optional[str], notebook_id: Optional[str]) -> Optional[str]:
    """与 kb_embedding 约定一致：返回该 notebook 的向量库根目录，供 RAG 使用。"""
    root = get_project_root()
    if not email:
        base = root / "outputs" / "kb_data" / "vector_store_main"
    else:
        safe_email = _sanitize_user_id(email)
        base = root / "outputs" / "kb_data" / safe_email
        if notebook_id:
            safe_nb = notebook_id.replace("/", "_").replace("\\", "_")[:128]
            base = base / safe_nb / "vector_store"
        else:
            base = base / "_shared" / "vector_store"
    return str(base) if base.exists() else None


def _resolve_chat_files(files: List[str]) -> List[str]:
    project_root = get_project_root()
    local_files: List[str] = []
    for f in files:
        clean_path = f.lstrip('/')
        p = project_root / clean_path
        if p.exists():
            local_files.append(str(p))
            log.info(f"[chat] ✓ Found file: {f} -> {p}")
            continue
        p_raw = Path(f)
        if p_raw.exists():
            local_files.append(str(p_raw))
            log.info(f"[chat] ✓ Found file (raw): {f} -> {p_raw}")
        else:
            log.warning(f"[chat] ✗ File not found: {f}")
    return local_files


def _resolve_vector_store_dir(email: Optional[str], notebook_id: Optional[str]) -> Optional[str]:
    vector_store_base_dir = None
    if email and notebook_id:
        try:
            project_root = get_project_root()
            email_dir = project_root / "outputs" / email.replace("@", "_at_")
            if email_dir.exists():
                for nb_dir in email_dir.iterdir():
                    if nb_dir.is_dir() and nb_dir.name.endswith(f"_{notebook_id}"):
                        vector_store_path = nb_dir / "vector_store"
                        if vector_store_path.exists():
                            vector_store_base_dir = str(vector_store_path)
                            log.info(f"[chat] Found vector store: {vector_store_base_dir}")
                            break
            if not vector_store_base_dir:
                log.warning(f"[chat] No vector_store found for email={email}, notebook_id={notebook_id}")
        except Exception as e:
            log.warning(f"[chat] Failed to search for vector store: {e}")

    if not vector_store_base_dir:
        vector_store_base_dir = _vector_store_base_dir(email, notebook_id)
        if vector_store_base_dir:
            log.info(f"[chat] Using legacy paths system, vector_store_base_dir: {vector_store_base_dir}")
        else:
            log.warning("[chat] vector_store_base_dir not found in either new or legacy system")
    return vector_store_base_dir


def _build_chat_request(
    files: List[str],
    query: str,
    history: List[Dict[str, str]],
    email: Optional[str],
    notebook_id: Optional[str],
    api_url: Optional[str],
    api_key: Optional[str],
    model: str,
) -> IntelligentQARequest:
    local_files = _resolve_chat_files(files)
    if not local_files:
        log.warning("[chat] No valid local files found, will rely on RAG only")

    if (api_url or "").strip() or (api_key or "").strip():
        log.warning("[chat] Ignoring frontend-supplied LLM config for notebook chat; using backend env config.")

    resolved_api_url, resolved_api_key = _require_llm_config(None, None)

    return IntelligentQARequest(
        file_ids=local_files,
        query=query,
        history=history,
        vector_store_base_dir=_resolve_vector_store_dir(email, notebook_id),
        chat_api_url=resolved_api_url,
        api_key=resolved_api_key,
        model=model,
    )


def _resolve_llm_api_url(explicit_api_url: Optional[str] = None) -> str:
    return (
        (explicit_api_url or "").strip()
        or str(getattr(settings, "LLM_API_URL", "") or "").strip()
        or str(os.getenv("LLM_API_URL", "")).strip()
        or str(os.getenv("DF_API_URL", "")).strip()
    )


def _resolve_llm_api_key(explicit_api_key: Optional[str] = None) -> str:
    return (
        (explicit_api_key or "").strip()
        or str(getattr(settings, "LLM_API_KEY", "") or "").strip()
        or str(os.getenv("LLM_API_KEY", "")).strip()
        or str(os.getenv("DF_API_KEY", "")).strip()
        or str(os.getenv("OPENAI_API_KEY", "")).strip()
    )


def _require_llm_config(
    explicit_api_url: Optional[str] = None,
    explicit_api_key: Optional[str] = None,
) -> tuple[str, str]:
    explicit_url = (explicit_api_url or "").strip()
    explicit_key = (explicit_api_key or "").strip()

    backend_api_url = _resolve_llm_api_url(None)
    backend_api_key = _resolve_llm_api_key(None)

    if explicit_url and explicit_key:
        resolved_api_url = explicit_url
        resolved_api_key = explicit_key
    elif explicit_url or explicit_key:
        log.warning(
            "[llm_config] Received incomplete explicit LLM config from request. "
            "Falling back to backend env pair to avoid mismatched url/key."
        )
        resolved_api_url = backend_api_url
        resolved_api_key = backend_api_key
    else:
        resolved_api_url = backend_api_url
        resolved_api_key = backend_api_key

    if not resolved_api_url or not resolved_api_key:
        raise HTTPException(
            status_code=400,
            detail="Missing LLM API configuration on the backend. Please set LLM_API_URL and LLM_API_KEY in fastapi_app/.env.",
        )

    return resolved_api_url, resolved_api_key


def _require_backend_llm_config() -> tuple[str, str]:
    return _require_llm_config(None, None)


def _resolve_search_provider(explicit_provider: Optional[str] = None) -> str:
    return ((explicit_provider or "").strip() or str(getattr(settings, "SEARCH_PROVIDER", "serper"))).lower()


def _resolve_search_api_key(provider: str, explicit_api_key: Optional[str] = None) -> str:
    if explicit_api_key and explicit_api_key.strip():
        return explicit_api_key.strip()

    provider = _resolve_search_provider(provider)
    if provider == "serper":
        return str(getattr(settings, "SERPER_API_KEY", "") or os.getenv("SERPER_API_KEY", "")).strip()
    if provider == "serpapi":
        return str(getattr(settings, "SERPAPI_KEY", "") or os.getenv("SERPAPI_KEY", "")).strip()
    if provider == "bocha":
        return str(getattr(settings, "BOCHA_API_KEY", "") or os.getenv("BOCHA_API_KEY", "")).strip()
    return ""


def _jsonl_line(payload: Dict[str, Any]) -> bytes:
    return (json.dumps(payload, ensure_ascii=False) + "\n").encode("utf-8")


def _require_workflow_available(name: str, *, feature_label: str, missing_detail: Optional[str] = None) -> None:
    available = list_workflows()
    if name in available:
        return

    detail = missing_detail or (
        f"{feature_label} 所需工作流 '{name}' 不可用。"
        f" 当前可用工作流: {list(available.keys())}"
    )
    raise HTTPException(status_code=503, detail=detail)


@router.post("/chat")
async def chat_with_kb(
    files: List[str] = Body(..., embed=True),
    query: str = Body(..., embed=True),
    history: List[Dict[str, str]] = Body([], embed=True),
    email: Optional[str] = Body(None, embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    api_url: Optional[str] = Body(None, embed=True),
    api_key: Optional[str] = Body(None, embed=True),
    model: str = Body(settings.KB_CHAT_MODEL, embed=True),
):
    """
    Intelligent QA Chat. 若传 email/notebook_id 且该 notebook 已建索引，会优先用 RAG 检索片段作为上下文。
    """
    log.info(f"[chat_with_kb] === Request received ===")
    log.info(f"[chat_with_kb] files (raw): {files}")
    log.info(f"[chat_with_kb] email (raw): {email}")
    log.info(f"[chat_with_kb] notebook_id (raw): {notebook_id}")
    log.info(f"[chat_with_kb] query length: {len(query)}")

    try:
        req = _build_chat_request(files, query, history, email, notebook_id, api_url, api_key, model)

        state = IntelligentQAState(request=req)
        
        # Run workflow via registry (统一使用 run_workflow)
        result_state = await run_workflow("intelligent_qa", state)
        
        # graph.ainvoke returns the final state dict or state object depending on implementation.
        # LangGraph usually returns dict. But our GenericGraphBuilder wrapper might return state.
        # GenericGraphBuilder compile returns a compiled graph.
        # Let's check typical usage. usually await graph.ainvoke(state) returns dict.
        
        answer = ""
        file_analyses = []
        source_mapping = {}
        source_preview_mapping = {}
        source_reference_mapping = {}

        if isinstance(result_state, dict):
            answer = result_state.get("answer", "")
            file_analyses = result_state.get("file_analyses", [])
            source_mapping = result_state.get("source_mapping", {})
            source_preview_mapping = result_state.get("source_preview_mapping", {})
            source_reference_mapping = result_state.get("source_reference_mapping", {})
        else:
            answer = getattr(result_state, "answer", "")
            file_analyses = getattr(result_state, "file_analyses", [])
            source_mapping = getattr(result_state, "source_mapping", {})
            source_preview_mapping = getattr(result_state, "source_preview_mapping", {})
            source_reference_mapping = getattr(result_state, "source_reference_mapping", {})

        # 将 source_mapping 的 int key 转为 str（JSON 要求）
        source_mapping_str = {str(k): v for k, v in source_mapping.items()} if source_mapping else {}
        source_preview_mapping_str = {str(k): v for k, v in source_preview_mapping.items()} if source_preview_mapping else {}
        source_reference_mapping_str = {str(k): v for k, v in source_reference_mapping.items()} if source_reference_mapping else {}

        return {
            "success": True,
            "answer": answer,
            "file_analyses": file_analyses,
            "source_mapping": source_mapping_str,
            "source_preview_mapping": source_preview_mapping_str,
            "source_reference_mapping": source_reference_mapping_str,
        }

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/chat/stream")
async def chat_with_kb_stream(
    files: List[str] = Body(..., embed=True),
    query: str = Body(..., embed=True),
    history: List[Dict[str, str]] = Body([], embed=True),
    email: Optional[str] = Body(None, embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    api_url: Optional[str] = Body(None, embed=True),
    api_key: Optional[str] = Body(None, embed=True),
    model: str = Body(settings.KB_CHAT_MODEL, embed=True),
):
    log.info("[chat_with_kb_stream] === Request received ===")

    async def event_generator():
        full_answer = ""
        try:
            req = _build_chat_request(files, query, history, email, notebook_id, api_url, api_key, model)
            state = IntelligentQAState(request=req)
            yield _jsonl_line({
                "type": "stage",
                "stage": "preparing",
                "message": "正在准备来源",
                "message_en": "Preparing sources",
            })

            yield _jsonl_line({
                "type": "stage",
                "stage": "analyzing",
                "message": "正在分析来源内容",
                "message_en": "Analyzing sources",
            })
            await prepare_parallel_file_analyses(state)

            yield _jsonl_line({
                "type": "stage",
                "stage": "retrieving",
                "message": "正在检索相关片段",
                "message_en": "Retrieving relevant chunks",
            })
            prompt = build_intelligent_qa_prompt(state)

            source_mapping_str = {str(k): v for k, v in (state.source_mapping or {}).items()}
            source_preview_mapping_str = {str(k): v for k, v in (state.source_preview_mapping or {}).items()}
            source_reference_mapping_str = {str(k): v for k, v in (state.source_reference_mapping or {}).items()}

            yield _jsonl_line({
                "type": "meta",
                "file_analyses": state.file_analyses,
                "source_mapping": source_mapping_str,
                "source_preview_mapping": source_preview_mapping_str,
                "source_reference_mapping": source_reference_mapping_str,
            })

            yield _jsonl_line({
                "type": "stage",
                "stage": "generating",
                "message": "正在生成回答",
                "message_en": "Generating answer",
            })

            client = AsyncOpenAI(
                api_key=req.api_key,
                base_url=req.chat_api_url,
            )

            stream = await client.chat.completions.create(
                model=req.model,
                messages=[
                    {"role": "system", "content": KbPromptAgentPrompts.system_prompt_for_kb_prompt_agent.strip()},
                    {"role": "user", "content": prompt},
                ],
                temperature=0.7,
                stream=True,
            )

            async for chunk in stream:
                delta = ""
                try:
                    delta = chunk.choices[0].delta.content or ""
                except Exception:
                    delta = ""
                if not delta:
                    continue
                full_answer += delta
                yield _jsonl_line({"type": "delta", "delta": delta})

            yield _jsonl_line({"type": "done", "answer": full_answer})
        except Exception as e:
            log.exception(f"[chat_with_kb_stream] failed: {e}")
            yield _jsonl_line({"type": "error", "message": str(e), "answer": full_answer})

    return StreamingResponse(event_generator(), media_type="application/x-ndjson")


# ---------- 1.1 对话记录：入库与读取 ----------
def _supabase_create_conversation(
    email: str,
    user_id: Optional[str],
    notebook_id: Optional[str],
    title: str = "新对话",
) -> Optional[Dict[str, Any]]:
    sb = get_supabase_admin_client()
    if not sb:
        return None
    try:
        ins = sb.table("kb_conversations").insert({
            "user_email": email,
            "user_id": user_id,
            "notebook_id": notebook_id,
            "title": title,
        }).execute()
        data = (ins.data or []) if hasattr(ins, "data") else []
        return data[0] if data else None
    except Exception as e:
        log.warning("supabase conversation create failed: %s", e)
        return None


@router.get("/conversations")
async def list_conversations_get(
    email: Optional[str] = None,
    user_id: Optional[str] = None,
    notebook_id: Optional[str] = None,
) -> Dict[str, Any]:
    """List conversations for user (by email or user_id), optionally filter by notebook_id."""
    sb = get_supabase_admin_client()
    if not sb:
        return {"success": True, "conversations": []}
    try:
        q = sb.table("kb_conversations").select("id,notebook_id,title,created_at,updated_at")
        if email:
            q = q.eq("user_email", email)
        if user_id:
            q = q.eq("user_id", user_id)
        if notebook_id:
            q = q.eq("notebook_id", notebook_id)
        r = q.order("updated_at", desc=True).limit(50).execute()
        rows = (r.data or []) if hasattr(r, "data") else []
        return {"success": True, "conversations": rows}
    except Exception as e:
        log.warning("list_conversations failed: %s", e)
        return {"success": True, "conversations": []}


@router.post("/conversations")
async def create_conversation(
    email: str = Body(..., embed=True),
    user_id: Optional[str] = Body(None, embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
) -> Dict[str, Any]:
    """Always create a new conversation for this user+notebook. Returns conversation id."""
    conv = _supabase_create_conversation(email, user_id, notebook_id)
    if conv:
        return {"success": True, "conversation_id": conv.get("id"), "conversation": conv}
    return {"success": False, "conversation_id": None}


@router.get("/conversations/{conversation_id}/messages")
async def get_conversation_messages(conversation_id: str) -> Dict[str, Any]:
    sb = get_supabase_admin_client()
    if not sb:
        return {"success": True, "messages": []}
    try:
        r = sb.table("kb_chat_messages").select("id,role,content,created_at").eq(
            "conversation_id", conversation_id
        ).order("created_at", desc=False).execute()
        rows = (r.data or []) if hasattr(r, "data") else []
        return {"success": True, "messages": rows}
    except Exception as e:
        log.warning("get_conversation_messages failed: %s", e)
        return {"success": True, "messages": []}


@router.post("/conversations/{conversation_id}/messages")
async def append_conversation_messages(
    conversation_id: str,
    messages: List[Dict[str, str]] = Body(..., embed=True),
) -> Dict[str, Any]:
    """Append messages (list of {role, content})."""
    sb = get_supabase_admin_client()
    if not sb:
        return {"success": False, "message": "Database not configured"}
    try:
        rows = [{"conversation_id": conversation_id, "role": m.get("role", "user"), "content": m.get("content", "")} for m in messages]
        sb.table("kb_chat_messages").insert(rows).execute()
        first_user_message = next((m.get("content", "").strip() for m in messages if m.get("role") == "user" and m.get("content", "").strip()), "")
        next_updated_at = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
        if first_user_message:
            conv = sb.table("kb_conversations").select("title").eq("id", conversation_id).limit(1).execute()
            conv_rows = (conv.data or []) if hasattr(conv, "data") else []
            current_title = str(conv_rows[0].get("title") or "").strip() if conv_rows else ""
            update_payload = {"updated_at": next_updated_at}
            if current_title in {"", "对话", "新对话"}:
                next_title = first_user_message[:28].strip()
                if next_title:
                    update_payload["title"] = next_title
            sb.table("kb_conversations").update(update_payload).eq("id", conversation_id).execute()
        else:
            sb.table("kb_conversations").update({"updated_at": next_updated_at}).eq("id", conversation_id).execute()
        return {"success": True}
    except Exception as e:
        log.warning("append_conversation_messages failed: %s", e)
        return {"success": False, "message": str(e)}


# ---------- 1.2 生成记录持久化：列表与写入 ----------
@router.get("/outputs-legacy")
async def list_outputs(
    email: Optional[str] = None,
    user_id: Optional[str] = None,
    notebook_id: Optional[str] = None,
    notebook_title: Optional[str] = None,
) -> Dict[str, Any]:
    """List generated outputs. Reads from JSON records first, falls back to filesystem scan."""
    from fastapi_app.kb_records import get_output_records

    em = email or user_id
    if not em:
        raise HTTPException(status_code=400, detail="email or user_id is required")

    # Try JSON records first (only if notebook_id is provided)
    if notebook_id:
        records = get_output_records(user_email=em, notebook_id=notebook_id)
        log.info(f"[list_outputs] Got {len(records)} records from JSON for {em}/{notebook_id}")
    else:
        records = []

    if records:
        files = []
        for r in records:
            files.append({
                "id": r.get("id") or f"output-{r['output_type']}-{int(r.get('created_at', 0))}",
                "output_type": r["output_type"],
                "file_name": r["file_name"],
                "title": r.get("title") or r["file_name"],
                "download_url": r["download_url"],
                "created_at": r.get("updated_at") or r.get("created_at"),
            })
        return {"success": True, "files": files}

    # Fallback to filesystem scan
    files: List[Dict[str, Any]] = []
    project_root = get_project_root()
    if not files and notebook_id:
        _FEATURE_EXT_MAP = {
            "ppt":     {".pdf", ".pptx"},
            "mindmap": {".mmd", ".mermaid"},
            "podcast": {".wav", ".mp3", ".m4a"},
            "drawio":  {".drawio"},
        }
        try:
            paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
            nb_root = paths.root
            if nb_root.exists():
                for feature, exts in _FEATURE_EXT_MAP.items():
                    feature_dir = nb_root / feature
                    if not feature_dir.exists():
                        continue
                    for ts_dir in feature_dir.iterdir():
                        if not ts_dir.is_dir():
                            continue
                        for f in ts_dir.iterdir():
                            if f.suffix.lower() in exts:
                                rel = str(f.relative_to(project_root))
                                files.append({
                                    "id": f"disk_{ts_dir.name}_{f.name}",
                                    "output_type": feature,
                                    "file_name": f.name,
                                    "download_url": _to_outputs_url(rel),
                                    "created_at": ts_dir.stat().st_mtime,
                                })
                                break
        except Exception as e:
            log.warning("list_outputs disk scan failed: %s", e)
    return {"success": True, "files": files}


def _extract_text_from_files(file_paths: List[str], max_chars: int = 50000) -> str:
    """从知识库文件列表中提取并合并文本，供 DrawIO 等使用。"""
    parts = []
    total = 0
    for f in file_paths:
        if total >= max_chars:
            break
        path = Path(f)
        if not path.exists():
            parts.append(f"[File not found: {f}]\n")
            total += len(parts[-1])
            continue
        suffix = path.suffix.lower()
        raw = ""
        try:
            if suffix == ".pdf":
                doc = fitz.open(f)
                raw = "\n".join(page.get_text() for page in doc)
                doc.close()
            elif suffix in [".docx", ".doc"]:
                try:
                    from docx import Document
                    doc = Document(f)
                    raw = "\n".join(p.text for p in doc.paragraphs)
                except Exception:
                    raw = "[Error: unsupported or missing python-docx]"
            elif suffix in [".pptx", ".ppt"]:
                try:
                    from pptx import Presentation
                    prs = Presentation(f)
                    raw = ""
                    for i, slide in enumerate(prs.slides):
                        raw += f"--- Slide {i+1} ---\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                raw += shape.text + "\n"
                except Exception:
                    raw = "[Error: unsupported or missing python-pptx]"
            else:
                with open(path, "r", encoding="utf-8", errors="replace") as fp:
                    raw = fp.read()
        except Exception as e:
            raw = f"[Parse Error: {e}]"
        chunk = (raw[: max_chars - total] + ("..." if len(raw) > max_chars - total else "")) if raw else ""
        parts.append(chunk)
        total += len(chunk)
    return "\n\n".join(parts)


def _save_output_record(
    email: str,
    user_id: Optional[str],
    notebook_id: Optional[str],
    output_type: str,
    file_name: str,
    file_path: str,
    result_path: str,
    download_url: str,
):
    sb = get_supabase_admin_client()
    if not sb:
        return
    try:
        sb.table("kb_output_records").insert({
            "user_email": email,
            "user_id": user_id,
            "notebook_id": notebook_id,
            "output_type": output_type,
            "file_name": file_name,
            "file_path": file_path,
            "result_path": result_path,
            "download_url": download_url,
        }).execute()
    except Exception as e:
        log.warning("_save_output_record failed: %s", e)


# 不做用户管理时使用的默认用户，数据从 outputs/local 取
DEFAULT_USER_ID = "local"
DEFAULT_EMAIL = "local"


@router.post("/fast-research")
async def fast_research(
    query: str = Body(..., embed=True),
    top_k: int = Body(10, embed=True),
    search_provider: Optional[str] = Body(None, embed=True),
    search_api_key: Optional[str] = Body(None, embed=True),
    search_engine: Optional[str] = Body(None, embed=True),
    google_cse_id: Optional[str] = Body(None, embed=True),
) -> Dict[str, Any]:
    """
    Fast Research: 用户输入查询，搜索引擎搜索，返回 top_k 条结果作为候选来源。
    支持：
    - serper：Google（环境变量 SERPER_API_KEY）
    - serpapi：Google / 百度，传 search_api_key + search_engine（google | baidu）
    - google_cse：传 search_api_key + google_cse_id
    - brave：传 search_api_key
    - bocha：博查 AI 网页搜索（https://api.bocha.cn），传 search_api_key（Bearer 鉴权）
    """
    top_k = max(1, min(20, top_k))
    resolved_provider = _resolve_search_provider(search_provider)
    resolved_search_api_key = _resolve_search_api_key(resolved_provider, search_api_key)
    sources = fast_research_search(
        query,
        top_k=top_k,
        search_provider=resolved_provider,
        search_api_key=resolved_search_api_key or None,
        search_engine=search_engine or "google",
        google_cse_id=google_cse_id,
    )
    return {
        "success": True,
        "query": query,
        "sources": sources,
    }


def _load_link_sources(nb_dir: Path) -> List[Dict[str, Any]]:
    path = nb_dir / LINK_SOURCES_FILENAME
    if not path.exists():
        return []
    try:
        raw = path.read_text(encoding="utf-8")
        data = json.loads(raw)
        return data if isinstance(data, list) else []
    except Exception:
        return []


def _save_link_sources(nb_dir: Path, items: List[Dict[str, Any]]) -> None:
    nb_dir.mkdir(parents=True, exist_ok=True)
    path = nb_dir / LINK_SOURCES_FILENAME
    path.write_text(json.dumps(items, ensure_ascii=False, indent=2), encoding="utf-8")


def _resolve_link_to_local_md(email: str, notebook_id: Optional[str], link_url: str) -> Optional[Path]:
    """
    若该 link 在「引入」时已抓取并存为本地 .md，则返回该 .md 的 Path，否则返回 None。
    这样生成时优先用本地已存文件，不再重新爬。
    """
    if not link_url or not (link_url.startswith("http://") or link_url.startswith("https://")):
        return None
    nb_dir = _notebook_dir(email or "default", notebook_id)
    link_list = _load_link_sources(nb_dir)
    link_url_stripped = link_url.strip()
    for item in link_list:
        if (item.get("link") or "").strip() != link_url_stripped:
            continue
        static_url = (item.get("static_url") or "").strip()
        if not static_url:
            return None
        try:
            p = _resolve_local_path(static_url)
            if p.exists() and p.is_file():
                return p
        except Exception:
            pass
    return None


@router.post("/import-link-sources")
async def import_link_sources(
    notebook_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    user_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    items: List[Dict[str, Any]] = Body(..., embed=True),
) -> Dict[str, Any]:
    """
    将 Fast Research 等返回的候选来源导入到当前笔记本。
    每个 URL 通过 httpx 抓取正文存为 .md，然后自动触发 embedding。
    New layout: outputs/{title}_{id}/sources/{stem}/
    """
    if not notebook_id or not email:
        raise HTTPException(status_code=400, detail="notebook_id and email are required")
    from fastapi_app.kb_records import add_source_record

    # New layout
    paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
    mgr = SourceManager(paths)

    # Legacy compat
    nb_dir = _notebook_dir(email, notebook_id)
    nb_dir.mkdir(parents=True, exist_ok=True)
    existing = _load_link_sources(nb_dir)
    seen_links = {x.get("link") for x in existing if x.get("link")}
    imported = 0
    saved_md_paths: List[str] = []

    for it in items:
        link = (it.get("link") or "").strip()
        if not link or link in seen_links:
            continue
        title = (it.get("title") or "").strip() or link
        snippet = (it.get("snippet") or "").strip()
        static_url = ""
        filename = ""
        try:
            text = await asyncio.to_thread(fetch_page_text, link)
            if not text or text.startswith("[抓取失败"):
                raise RuntimeError(text or "empty response")

            # Import into new layout
            source_info = await mgr.import_url(link, text, title[:80])
            project_root = get_project_root()
            rel = source_info.original_path.relative_to(project_root)
            static_url = "/" + rel.as_posix()
            filename = source_info.original_path.name
            saved_md_paths.append(str(source_info.original_path))

            # Legacy compat: also save to old path
            legacy_filename = _sanitize_md_filename(title[:80], "link")
            legacy_path = nb_dir / legacy_filename
            if not legacy_path.exists():
                try:
                    shutil.copy2(str(source_info.original_path), str(legacy_path))
                except Exception:
                    pass

            try:
                add_source_record(
                    user_email=email,
                    notebook_id=notebook_id,
                    file_name=source_info.original_path.name,
                    file_path=str(source_info.original_path),
                    static_url=static_url,
                    file_size=source_info.original_path.stat().st_size,
                    file_type="text/markdown",
                )
            except Exception as record_err:
                log.warning("[import-link-sources] failed to write source record for %s: %s", link[:60], record_err)

            log.info("[import-link-sources] 已抓取并保存: %s -> %s", link[:60], filename)
        except Exception as e:
            log.warning("[import-link-sources] 抓取失败 %s: %s", link[:60], e)

        existing.append({
            "id": f"link-{int(time.time() * 1000)}-{imported}",
            "title": title[:500],
            "link": link,
            "snippet": snippet[:2000],
            "static_url": static_url,
            "filename": filename,
        })
        seen_links.add(link)
        imported += 1
    _save_link_sources(nb_dir, existing)

    # Auto-embed saved .md files using new paths
    embedded = 0
    if saved_md_paths:
        try:
            vector_base = str(paths.vector_store_dir)
            file_list = [{"path": p} for p in saved_md_paths]
            await process_knowledge_base_files(
                file_list=file_list,
                base_dir=vector_base,
            )
            embedded = len(saved_md_paths)
            log.info("[import-link-sources] embedding 完成, %d 个文件", embedded)
        except Exception as e:
            log.warning("[import-link-sources] embedding 失败: %s", e)

    return {"success": True, "imported": imported, "embedded": embedded}


@router.post("/generate-ppt")
async def generate_ppt_from_kb(
    file_path: Optional[str] = Body(None, embed=True),
    file_paths: Optional[List[str]] = Body(None, embed=True),
    image_paths: Optional[List[str]] = Body(None, embed=True),
    image_items: Optional[List[Dict[str, Any]]] = Body(None, embed=True),
    query: Optional[str] = Body("", embed=True),
    need_embedding: bool = Body(False, embed=True),
    search_top_k: int = Body(8, embed=True),
    user_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    api_url: Optional[str] = Body(None, embed=True),
    api_key: Optional[str] = Body(None, embed=True),
    style: str = Body("modern", embed=True),
    language: str = Body("zh", embed=True),
    page_count: int = Body(10, embed=True),
    model: str = Body("deepseek-v3.2", embed=True),
    gen_fig_model: str = Body(settings.IMAGE_GEN_MODEL or settings.PAPER2PPT_IMAGE_GEN_MODEL, embed=True),
):
    """
    Generate PPT from knowledge base file. Outputs under user/notebook dir.
    """
    try:
        _require_workflow_available(
            "kb_page_content",
            feature_label="PPT 生成",
        )
        _require_workflow_available(
            "paper2ppt_parallel_consistent_style",
            feature_label="PPT 生成",
            missing_detail=(
                "PPT 生成功能当前不可用：后端未加载工作流 'paper2ppt_parallel_consistent_style'。"
                " 通常是缺少可选依赖 'paddle' / 'paddleocr' 导致启动时跳过了该工作流。"
            ),
        )

        api_url, api_key = _require_llm_config(api_url, api_key)
        # 兼容前端传 file_paths 为数组或单个字符串；保证多选时每项一个来源
        if file_paths is not None:
            raw_list = file_paths if isinstance(file_paths, list) else [file_paths] if file_paths else []
        else:
            raw_list = [file_path] if file_path else []
        input_paths = [x for x in raw_list if x]
        log.info("[generate-ppt] 收到 file_paths 数量: %s", len(input_paths))

        if not input_paths:
            raise HTTPException(status_code=400, detail="No input files provided")
        if not isinstance(page_count, int) or page_count < 1 or page_count > 50:
            raise HTTPException(status_code=400, detail="page_count must be an integer between 1 and 50")
        log.info("[generate-ppt] 收到 page_count=%s", page_count)

        # 区分本地文件与网页 URL（前端会传 type=link 的 url 为 http(s) 链接）
        url_sources: List[str] = []
        path_sources: List[Path] = []
        user_image_items: List[Dict[str, Any]] = []
        seen_resolved: set = set()
        for p in input_paths:
            ps = (p or "").strip()
            if ps.startswith("http://") or ps.startswith("https://"):
                url_sources.append(ps)
                continue
            local_path = _resolve_local_path(p)
            if not local_path.exists():
                raise HTTPException(status_code=404, detail=f"File not found: {p}")
            ext = local_path.suffix.lower()
            if ext in IMAGE_EXTENSIONS:
                user_image_items.append({"path": str(local_path), "description": ""})
                continue
            if ext in {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".md", ".markdown"}:
                key = str(local_path.resolve())
                if key not in seen_resolved:
                    seen_resolved.add(key)
                    path_sources.append(local_path)
            else:
                raise HTTPException(status_code=400, detail=f"Unsupported file type for PPT: {local_path.name}")

        if not path_sources and not url_sources:
            raise HTTPException(status_code=400, detail="At least one document or web source is required for PPT generation")

        ts = int(time.time())
        project_root = get_project_root()
        # New layout: outputs/{title}_{id}/ppt/{ts}/
        if notebook_id:
            nb_paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
            output_dir = nb_paths.feature_output_dir("ppt", ts)
        else:
            output_dir = _outputs_dir(email, notebook_id, f"{ts}_ppt")
        output_dir.mkdir(parents=True, exist_ok=True)

        md_exts = {".md", ".markdown"}
        pdf_like_exts = {".pdf", ".pptx", ".ppt", ".docx", ".doc"}
        doc_paths = path_sources
        md_paths = [p for p in doc_paths if p.suffix.lower() in md_exts]
        pdf_like_paths = [p for p in doc_paths if p.suffix.lower() in pdf_like_exts]

        use_text_input = bool(md_paths) or bool(url_sources)
        combined_text = ""
        local_file_path = None
        pdf_paths_for_outline: List[Path] = []

        if use_text_input:
            # 按 input_paths 顺序：先本地文件再 URL，生成「来源1」「来源2」…（含网页抓取）
            text_parts: List[str] = []
            idx = 0
            for p in input_paths:
                ps = (p or "").strip()
                if ps.startswith("http://") or ps.startswith("https://"):
                    idx += 1
                    content = None
                    local_md = _resolve_link_to_local_md(email, notebook_id, ps)
                    if local_md is not None:
                        try:
                            content = local_md.read_text(encoding="utf-8", errors="replace")
                            log.info("[generate-ppt] 网页来源使用已存 .md: %s", local_md.name)
                        except Exception as e:
                            log.warning("[generate-ppt] 读取已存 .md 失败 %s: %s", local_md, e)
                    if not (content or "").strip():
                        try:
                            content = fetch_page_text(ps, max_chars=100000)
                            if content:
                                log.info("[generate-ppt] 网页来源 %s 抓取成功，长度=%s", idx, len(content))
                        except Exception as e:
                            log.warning("[generate-ppt] 网页来源抓取失败 %s: %s", ps[:80], e)
                    if (content or "").strip():
                        text_parts.append(f"来源{len(text_parts) + 1}:\n{content.strip()}")
                    continue
                local_path = _resolve_local_path(p)
                if not local_path.exists():
                    continue
                ext = local_path.suffix.lower()
                if ext not in {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".md", ".markdown"}:
                    continue
                idx += 1
                try:
                    if ext in md_exts:
                        content = local_path.read_text(encoding="utf-8")
                    elif ext == ".pdf":
                        # 优先从 MinerU 缓存读取高质量 markdown
                        content = _read_mineru_md_if_cached(local_path, email, notebook_id, notebook_title=notebook_title)
                        if not content:
                            content = _extract_text_from_files([str(local_path)])
                    elif ext in pdf_like_exts:
                        content = _extract_text_from_files([str(local_path)])
                    else:
                        content = ""
                    if (content or "").strip():
                        text_parts.append(f"来源{len(text_parts) + 1}:\n{content.strip()}")
                except Exception as e:
                    log.warning("read doc %s (来源%s): %s", local_path.name, len(text_parts) + 1, e)
            combined_text = "\n\n".join(text_parts).strip()
            log.info("[generate-ppt] 共 %s 个来源（本地 %s + 网页 %s），TEXT 块数: %s", len(path_sources) + len(url_sources), len(path_sources), len(url_sources), len(text_parts))
            if not combined_text:
                raise HTTPException(status_code=400, detail="No text content could be read from the selected sources")
        else:
            # 仅 PDF/PPTX/DOCX：转 PDF 后合并
            local_pdf_paths: List[Path] = []
            convert_dir = output_dir / "input"
            convert_dir.mkdir(parents=True, exist_ok=True)
            for p in pdf_like_paths:
                ext = p.suffix.lower()
                if ext == ".pdf":
                    local_pdf_paths.append(p)
                elif ext in {".pptx", ".ppt", ".docx", ".doc"}:
                    local_pdf_paths.append(_convert_to_pdf(p, convert_dir))
                else:
                    raise HTTPException(status_code=400, detail=f"Unsupported file type for PPT: {p.name}")

            pdf_paths_for_outline = local_pdf_paths
            if len(local_pdf_paths) > 1:
                merge_dir = output_dir / "input"
                merged_pdf = merge_dir / "merged.pdf"
                local_file_path = _merge_pdfs(local_pdf_paths, merged_pdf)
            else:
                local_file_path = local_pdf_paths[0]

        # Normalize image items (optional)
        resolved_image_items: List[Dict[str, Any]] = []
        for item in image_items or []:
            raw_path = item.get("path") or item.get("url") or ""
            if not raw_path:
                continue
            img_path = _resolve_local_path(str(raw_path))
            if img_path.exists() and img_path.suffix.lower() in IMAGE_EXTENSIONS:
                resolved_image_items.append({
                    "path": str(img_path),
                    "description": item.get("description") or item.get("desc") or ""
                })

        for img in image_paths or []:
            img_path = _resolve_local_path(img)
            if img_path.exists() and img_path.suffix.lower() in IMAGE_EXTENSIONS:
                resolved_image_items.append({
                    "path": str(img_path),
                    "description": ""
                })

        resolved_image_items.extend(user_image_items)

        # Embedding + retrieval (optional): use notebook-scoped vector store，入库只用本地 embedding，MinerU 输出到 kb_mineru
        retrieval_text = ""
        if need_embedding:
            base_dir = _notebook_dir(email, notebook_id) / "vector_store"
            embed_api_url = api_url
            if "/embeddings" not in embed_api_url:
                embed_api_url = embed_api_url.rstrip("/") + "/embeddings"
            project_root = get_project_root()
            safe_nb = (notebook_id or "_shared").replace("/", "_").replace("\\", "_")[:128]
            safe_email = _sanitize_user_id(email) if email else "default"
            mineru_output_base = project_root / "outputs" / "kb_mineru" / safe_email / safe_nb
            mineru_output_base.mkdir(parents=True, exist_ok=True)

            files_for_embed = [{"path": str(p), "description": ""} for p in doc_paths]
            manifest = await process_knowledge_base_files(
                files_for_embed,
                base_dir=str(base_dir),
                api_url=None,
                api_key=api_key,
                model_name=None,
                multimodal_model=None,
                mineru_output_base=str(mineru_output_base),
            )

            manager = VectorStoreManager(
                base_dir=str(base_dir),
                api_key=api_key,
            )

            def _match_file_ids(m: Dict[str, Any], paths: List[Path]) -> List[str]:
                ids: List[str] = []
                target = {str(p.resolve()) for p in paths}
                for f in m.get("files", []):
                    try:
                        if str(Path(f.get("original_path", "")).resolve()) in target:
                            if f.get("id"):
                                ids.append(f["id"])
                    except Exception:
                        continue
                return ids

            file_ids = _match_file_ids(manifest or manager.manifest or {}, doc_paths)
            if query and file_ids:
                results = manager.search(query=query, top_k=search_top_k, file_ids=file_ids)
                retrieval_text = "\n\n".join([r.get("content", "") for r in results if r.get("content")])

        # Prepare request（支持 PDF 或 TEXT：.md 及混合时用 TEXT）
        ppt_req = Paper2PPTRequest(
            input_type="TEXT" if use_text_input else "PDF",
            input_content=combined_text if use_text_input else str(local_file_path),
            email=email,
            chat_api_url=api_url,
            chat_api_key=api_key,
            api_key=api_key,
            style=style,
            language=language,
            page_count=page_count,
            model=model,
            gen_fig_model=gen_fig_model,
            aspect_ratio="16:9",
            use_long_paper=False
        )
        log.info("[generate-ppt] ppt_req.page_count=%s（将传入 outline 生成）", ppt_req.page_count)

        # 复用 embedding 入库时已有的 MinerU 解析结果，避免重复跑 MinerU
        if not use_text_input and pdf_paths_for_outline:
            n_reused = _reuse_mineru_cache(pdf_paths_for_outline, output_dir, email, notebook_id, notebook_title=notebook_title)
            log.info("[generate-ppt] MinerU 缓存复用: %s/%s 个 PDF", n_reused, len(pdf_paths_for_outline))

        # Step 1: 生成大纲（kb_page_content 内含 LLM outline_agent，无人工确认）
        log.info("[generate-ppt] Step 1: 运行 kb_page_content，由 LLM 生成大纲 (outline)")
        state_pc = _init_state_from_request(ppt_req, result_path=output_dir)
        state_pc.kb_query = query or ""
        state_pc.kb_retrieval_text = retrieval_text
        state_pc.kb_user_images = resolved_image_items
        # 多 PDF 时按「来源1:\n...\n\n来源2:\n...」拼入，供 outline 使用
        if not use_text_input and len(pdf_paths_for_outline) > 1:
            multi_parts = []
            for i, p in enumerate(pdf_paths_for_outline):
                # 优先从 MinerU 缓存读取高质量 markdown
                part = _read_mineru_md_if_cached(p, email, notebook_id, notebook_title=notebook_title)
                if not part:
                    part = _extract_text_from_files([str(p)])
                if part.strip():
                    multi_parts.append(f"来源{i + 1}:\n{part}")
            if multi_parts:
                state_pc.kb_multi_source_text = "\n\n".join(multi_parts)
        state_pc_result = await run_workflow("kb_page_content", state_pc)
        if isinstance(state_pc_result, dict):
            for k, v in state_pc_result.items():
                setattr(state_pc, k, v)
        else:
            state_pc = state_pc_result
        pagecontent = getattr(state_pc, "pagecontent", []) or []
        log.info("[generate-ppt] Step 1 完成: 大纲已生成，共 %s 页", len(pagecontent))
        if not pagecontent:
            raise HTTPException(status_code=500, detail="大纲生成结果为空，请检查输入文档或重试")

        # Step 2: 按大纲生图并导出 PDF/PPTX（与 Paper2Any 一致使用 paper2ppt_parallel_consistent_style）
        state_pc.pagecontent = pagecontent
        log.info("[generate-ppt] Step 2: 运行 paper2ppt_parallel_consistent_style 生图")
        state_pp = await run_workflow("paper2ppt_parallel_consistent_style", state_pc)

        # Extract output paths (workflow may set ppt_pdf_path / ppt_pptx_path)
        if isinstance(state_pp, dict):
            pdf_path = state_pp.get("ppt_pdf_path") or ""
            pptx_path = state_pp.get("ppt_pptx_path") or ""
        else:
            pdf_path = getattr(state_pp, "ppt_pdf_path", None) or ""
            pptx_path = getattr(state_pp, "ppt_pptx_path", None) or ""
        # 若 workflow 未写回路径，则按约定路径回退：output_dir 下 paper2ppt.pdf / paper2ppt_editable.pptx
        if not pdf_path:
            fallback_pdf = Path(output_dir) / "paper2ppt.pdf"
            if fallback_pdf.exists():
                pdf_path = str(fallback_pdf)
        if not pptx_path:
            fallback_pptx = Path(output_dir) / "paper2ppt_editable.pptx"
            if fallback_pptx.exists():
                pptx_path = str(fallback_pptx)

        pdf_url = _to_outputs_url(pdf_path) if pdf_path else ""
        pptx_url = _to_outputs_url(pptx_path) if pptx_path else ""
        # 下载链接优先 PDF（可预览），其次 PPTX
        download_url = pdf_url or pptx_url
        _save_output_record(
            email=email,
            user_id=user_id,
            notebook_id=notebook_id,
            output_type="ppt",
            file_name="paper2ppt.pdf",
            file_path=pdf_path or "",
            result_path=str(output_dir),
            download_url=download_url,
        )

        from fastapi_app.kb_records import add_output_record
        try:
            add_output_record(
                user_email=email,
                notebook_id=notebook_id,
                output_type="ppt",
                file_name="paper2ppt.pdf",
                download_url=download_url
            )
        except Exception as e:
            log.warning("[generate-ppt] failed to write JSON record: %s", e)

        return {
            "success": True,
            "result_path": str(output_dir),
            "pdf_path": pdf_url,
            "pptx_path": pptx_url,
            "download_url": download_url,
            "output_file_id": f"kb_ppt_{ts}",
        }

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/generate-deep-research-report")
async def generate_deep_research_report(
    topic: str = Body(..., embed=True),
    user_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    api_url: Optional[str] = Body(None, embed=True),
    api_key: Optional[str] = Body(None, embed=True),
    language: str = Body("zh", embed=True),
    style: str = Body("modern", embed=True),
    page_count: int = Body(10, embed=True),
    model: str = Body("deepseek-v3.2", embed=True),
    gen_fig_model: str = Body(settings.IMAGE_GEN_MODEL or settings.PAPER2PPT_IMAGE_GEN_MODEL, embed=True),
    add_as_source: bool = Body(True, embed=True),
    search_provider: Optional[str] = Body(None, embed=True),
    search_api_key: Optional[str] = Body(None, embed=True),
    search_engine: Optional[str] = Body(None, embed=True),
    search_top_k: int = Body(10, embed=True),
    # 新增：DeepResearch 完整模式配置
    use_full_deep_research: bool = Body(True, embed=True),  # 默认使用完整的阿里DeepResearch
    max_iterations: int = Body(50, embed=True),  # DeepResearch最大迭代次数
    serper_api_key: Optional[str] = Body(None, embed=True),  # Serper API密钥
    jina_api_key: Optional[str] = Body(None, embed=True),  # Jina API密钥
) -> Dict[str, Any]:
    """
    Deep Research 报告生成（默认使用完整版阿里DeepResearch）：
    - use_full_deep_research=True: 完整版（阿里DeepResearch多轮ReAct推理，深度）【默认】
    - use_full_deep_research=False: 简化版（搜索 + LLM总结，快速）
    """
    try:
        api_url, api_key = _require_llm_config(api_url, api_key)
        if not isinstance(page_count, int) or page_count < 1 or page_count > 50:
            raise HTTPException(status_code=400, detail="page_count must be an integer between 1 and 50")
        ts = int(time.time())
        project_root = get_project_root()

        # New layout: outputs/{title}_{id}/deep_research/{ts}/
        if notebook_id:
            dr_paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
            output_dir = dr_paths.feature_output_dir("deep_research", ts)
        else:
            output_dir = _outputs_dir(email, notebook_id, f"{ts}_deep_research")
        output_dir.mkdir(parents=True, exist_ok=True)

        topic = topic.strip()

        # ============================================================================
        # 模式选择：完整DeepResearch vs 简化版
        # ============================================================================

        if use_full_deep_research:
            # 使用完整的阿里DeepResearch（多轮ReAct推理）
            log.info("[generate-deep-research-report] 使用完整DeepResearch模式: topic=%r, max_iterations=%s", topic[:150], max_iterations)

            # 如果没有传递 serper_api_key，尝试使用 search_api_key 作为回退
            final_serper_key = serper_api_key or search_api_key or _resolve_search_api_key("serper", None)

            log.info("[generate-deep-research-report] API配置: serper_api_key=%s, search_api_key=%s, final_serper_key=%s",
                     "***" if serper_api_key else "None",
                     "***" if search_api_key else "None",
                     "***" if final_serper_key else "None")

            # 运行完整DeepResearch（直接传递参数，不依赖环境变量）
            from fastapi_app.services.deep_research_integration import DeepResearchIntegration

            integration = DeepResearchIntegration(
                model_name=model,
                api_base=api_url,
                api_key=api_key,
                max_iterations=max_iterations,
                serper_key=final_serper_key,
                jina_keys=jina_api_key,
            )
            result = await integration.run_research(
                query=topic,
                max_iterations=max_iterations
            )

            if not result["success"]:
                raise HTTPException(status_code=500, detail=result.get("error", "DeepResearch failed"))

            # 格式化为Markdown
            report = integration.format_result_as_markdown(result)
            report_title = f"DeepResearch: {topic[:50]}"

            log.info("[generate-deep-research-report] 完整DeepResearch完成: iterations=%s, sources=%s",
                     result.get("iterations", 0), len(result.get("sources", [])))

        else:
            # 使用简化版（搜索 + LLM总结）
            search_top_k = max(1, min(20, search_top_k))
            log.info(
                "[generate-deep-research-report] 使用简化版模式: topic=%r, search_top_k=%s, provider=%s, model=%s",
                topic[:150], search_top_k, search_provider, model,
            )

            # 1) 搜索：用 topic 做 Fast Research，拿到 top_k 条结果
            sources = fast_research_search(
                topic,
                top_k=search_top_k,
                search_provider=_resolve_search_provider(search_provider),
                search_api_key=(
                    search_api_key
                    or serper_api_key
                    or _resolve_search_api_key(search_provider or "serper", None)
                ),
                search_engine=search_engine or "google",
            )
            log.info("[generate-deep-research-report] search 完成: 共 %s 条来源", len(sources))

            search_context = ""
            if sources:
                search_context = "\n\n".join(
                    f"[{i+1}] 标题: {s.get('title', '')}\n链接: {s.get('link', '')}\n摘要: {s.get('snippet', '')}"
                    for i, s in enumerate(sources)
                )
                log.info("[generate-deep-research-report] search_context 拼接完成: len=%s", len(search_context))
            else:
                log.warning("[generate-deep-research-report] no search results, LLM will generate from topic only")

            # 2) LLM：根据 topic + search_context 生成一篇长报告（返回标题 + 正文）
            report_title, report = generate_report_from_search(
                topic=topic,
                search_context=search_context,
                api_url=api_url,
                api_key=api_key,
                model=model,
                language=language,
            )
            if not (report or "").strip():
                raise HTTPException(status_code=500, detail="LLM did not return report content")
            log.info("[generate-deep-research-report] 简化版报告生成完成: title=%r, report_len=%s", report_title, len(report))

        # 3) 来源名：固定前缀 [report] + LLM 给的标题，保存为 .md
        safe_title = re.sub(r'[/\\:*?"<>|]', "", (report_title or "").strip()) or "report"
        safe_title = safe_title[:50].strip()
        file_name = f"[report] {safe_title}_{ts}.md"
        report_path = output_dir / file_name
        log.info("[generate-deep-research-report] 开始写入 Markdown: %s", report_path)
        report_path.write_text(report, encoding="utf-8")
        if not report_path.exists():
            raise HTTPException(status_code=500, detail="Deep research report file was not written")

        report_url = _to_outputs_url(str(report_path))
        log.info("[generate-deep-research-report] 报告已保存: %s, add_as_source=%s, notebook_id=%s", report_path, add_as_source, notebook_id)

        if add_as_source and notebook_id:
            nb_dir = _notebook_dir(email, notebook_id)
            nb_dir.mkdir(parents=True, exist_ok=True)
            dest = nb_dir / file_name
            shutil.copy2(str(report_path), dest)
            try:
                rel = dest.relative_to(project_root)
                source_static_url = "/" + rel.as_posix().replace("@", "%40")
            except ValueError:
                source_static_url = report_url
            _save_output_record(
                email=email,
                user_id=user_id,
                notebook_id=notebook_id,
                output_type="report",
                file_name=file_name,
                file_path=str(dest),
                result_path=str(output_dir),
                download_url=report_url,
            )

            from fastapi_app.kb_records import add_output_record
            try:
                add_output_record(
                    user_email=email,
                    notebook_id=notebook_id,
                    output_type="report",
                    file_name=file_name,
                    download_url=report_url
                )
            except Exception as e:
                log.warning("[generate-deep-research-report] failed to write JSON record: %s", e)

            stat = dest.stat()
            added_file = {
                "id": f"file-{file_name}-{stat.st_mtime_ns}",
                "name": file_name,
                "url": source_static_url,
                "static_url": source_static_url,
                "file_size": stat.st_size,
                "file_type": "text/markdown",
            }
            log.info("[generate-deep-research-report] 完成: 已加入来源, file_name=%s", file_name)
            return {
                "success": True,
                "pdf_path": report_url,
                "pdf_url": report_url,
                "file_name": file_name,
                "source_static_url": source_static_url,
                "added_as_source": True,
                "added_file": added_file,
            }

        _save_output_record(
            email=email,
            user_id=user_id,
            notebook_id=notebook_id,
            output_type="report",
            file_name=file_name,
            file_path=str(report_path),
            result_path=str(output_dir),
            download_url=report_url,
        )

        from fastapi_app.kb_records import add_output_record
        try:
            add_output_record(
                user_email=email,
                notebook_id=notebook_id,
                output_type="report",
                file_name=file_name,
                download_url=report_url
            )
        except Exception as e:
            log.warning("[generate-deep-research-report] failed to write JSON record: %s", e)

        log.info("[generate-deep-research-report] 完成: 未加入来源, file_name=%s", file_name)
        return {
            "success": True,
            "pdf_path": report_url,
            "pdf_url": report_url,
            "file_name": file_name,
            "added_as_source": False,
        }
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/generate-podcast")
async def generate_podcast_from_kb(
    file_paths: List[str] = Body(..., embed=True),
    user_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    api_url: Optional[str] = Body(None, embed=True),
    api_key: Optional[str] = Body(None, embed=True),
    model: str = Body("deepseek-v3.2", embed=True),
    tts_model: str = Body("gemini-2.5-pro-preview-tts", embed=True),
    voice_name: str = Body("Cherry", embed=True),
    voice_name_b: str = Body("Chelsie", embed=True),
    podcast_mode: str = Body("monologue", embed=True),
    language: str = Body("zh", embed=True),
):
    """
    从知识库生成播客。支持本地文件与「搜索引入」的 URL：URL 优先用已存 .md，否则抓取后写临时 .md 再参与生成。
    """
    try:
        model = str(_unwrap_fastapi_body_default(model, settings.KB_CHAT_MODEL) or settings.KB_CHAT_MODEL)
        tts_model = str(_unwrap_fastapi_body_default(tts_model, settings.TTS_MODEL) or settings.TTS_MODEL)
        voice_name = str(_unwrap_fastapi_body_default(voice_name, "Cherry") or "Cherry")
        voice_name_b = str(_unwrap_fastapi_body_default(voice_name_b, "Chelsie") or "Chelsie")
        podcast_mode = str(_unwrap_fastapi_body_default(podcast_mode, "monologue") or "monologue")
        language = str(_unwrap_fastapi_body_default(language, "zh") or "zh")
        api_url, api_key = _require_llm_config(api_url, api_key)
        # Validate TTS mode restrictions
        if podcast_mode == "dialog":
            tts_lower = tts_model.lower()
            if "qwen" in tts_lower:
                raise HTTPException(status_code=400, detail="qwen-tts 仅支持单人播客模式")
            if "gemini" in tts_lower and "tts" in tts_lower and "apiyi" not in api_url.lower():
                raise HTTPException(status_code=400, detail="gemini-2.5-flash-tts 双人模式需要使用 apiyi 平台")

        ts = int(time.time())
        # New layout: outputs/{title}_{id}/podcast/{ts}/
        if notebook_id:
            paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
            output_dir = paths.feature_output_dir("podcast", ts)
        else:
            output_dir = _outputs_dir(email, notebook_id, f"{ts}_podcast")
        output_dir.mkdir(parents=True, exist_ok=True)
        project_root = get_project_root()

        if not file_paths:
            raise HTTPException(status_code=400, detail="No valid files provided")

        local_paths: List[Path] = []
        for f in (file_paths or []):
            ps = (f or "").strip()
            if ps.startswith("http://") or ps.startswith("https://"):
                content = None
                local_md = _resolve_link_to_local_md(email, notebook_id, ps)
                if local_md is not None:
                    try:
                        content = local_md.read_text(encoding="utf-8", errors="replace")
                        log.info("[generate-podcast] 网页来源使用已存 .md: %s", local_md.name)
                    except Exception as e:
                        log.warning("[generate-podcast] 读取已存 .md 失败: %s", e)
                if not (content or "").strip():
                    try:
                        content = fetch_page_text(ps, max_chars=100000)
                    except Exception as e:
                        log.warning("[generate-podcast] 抓取 URL 失败 %s: %s", ps[:60], e)
                        content = ""
                if (content or "").strip():
                    link_dir = output_dir / "input"
                    link_dir.mkdir(parents=True, exist_ok=True)
                    tmp_md = link_dir / f"link_{len(local_paths)}.md"
                    tmp_md.write_text(content.strip(), encoding="utf-8")
                    local_paths.append(tmp_md)
            else:
                local_path = _resolve_local_path(ps)
                if not local_path.exists() or not local_path.is_file():
                    raise HTTPException(status_code=404, detail=f"File not found: {ps}")
                local_paths.append(local_path)

        # 过滤不支持的文件类型（例如图片），只保留可转文本的文档（含 .md 报告）
        supported_exts = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".md", ".markdown"}
        filtered_paths: List[Path] = []
        ignored_paths: List[Path] = []
        for p in local_paths:
            if p.suffix.lower() in supported_exts:
                filtered_paths.append(p)
            else:
                ignored_paths.append(p)

        if not filtered_paths:
            raise HTTPException(status_code=400, detail="No supported document files for podcast (support: PDF, Word, PPT, MD)")

        if ignored_paths:
            log.warning(
                "[kb_podcast] ignore unsupported files: "
                + ", ".join([p.name for p in ignored_paths])
            )

        # If multiple files, merge into a single PDF (doc/ppt will be converted); 仅 .md 时拼成一份 .md
        if len(filtered_paths) > 1:
            merge_dir = output_dir / "input"
            merge_dir.mkdir(parents=True, exist_ok=True)

            pdf_paths: List[Path] = []
            md_paths: List[Path] = []
            for p in filtered_paths:
                ext = p.suffix.lower()
                if ext == ".pdf":
                    pdf_paths.append(p)
                elif ext in {".docx", ".doc", ".pptx", ".ppt"}:
                    pdf_paths.append(_convert_to_pdf(p, merge_dir))
                elif ext in {".md", ".markdown"}:
                    md_paths.append(p)

            if pdf_paths:
                merged_pdf = merge_dir / "merged.pdf"
                local_file_paths = [str(_merge_pdfs(pdf_paths, merged_pdf))]
            elif md_paths:
                merged_md = merge_dir / "merged.md"
                parts = [p.read_text(encoding="utf-8", errors="replace") for p in md_paths]
                merged_md.write_text("\n\n".join(parts), encoding="utf-8")
                local_file_paths = [str(merged_md)]
            else:
                raise HTTPException(status_code=400, detail="No supported document files for podcast")
        else:
            local_file_paths = [str(filtered_paths[0])]

        # Get vector store base directory
        vector_store_base_dir = _vector_store_base_dir(email, notebook_id)

        # Prepare request
        podcast_req = KBPodcastRequest(
            file_ids=local_file_paths,
            vector_store_base_dir=vector_store_base_dir,
            chat_api_url=api_url,
            api_key=api_key,
            model=model,
            tts_model=tts_model,
            voice_name=voice_name,
            voice_name_b=voice_name_b,
            podcast_mode=podcast_mode,
            language=language
        )
        podcast_req.email = email

        state = KBPodcastState(request=podcast_req, result_path=str(output_dir))

        # Run workflow via registry (统一使用 run_workflow)
        result_state = await run_workflow("kb_podcast", state)

        # Extract results
        audio_path = ""
        script_path = ""
        result_path = ""

        if isinstance(result_state, dict):
            audio_path = result_state.get("audio_path", "")
            result_path = result_state.get("result_path", "")
        else:
            audio_path = getattr(result_state, "audio_path", "")
            result_path = getattr(result_state, "result_path", "")

        if result_path:
            script_path = str(Path(result_path) / "script.txt")

        audio_error = ""
        if not audio_path:
            audio_error = "No audio path returned from workflow"
        elif isinstance(audio_path, str) and audio_path.startswith("["):
            audio_error = audio_path
        else:
            audio_file = Path(audio_path)
            if not audio_file.is_absolute():
                audio_file = (get_project_root() / audio_file).resolve()
            if not audio_file.exists():
                audio_error = f"Audio file not found: {audio_file}"

        if audio_error:
            raise HTTPException(status_code=500, detail=audio_error)

        audio_url = _to_outputs_url(audio_path) if audio_path else ""
        script_url = _to_outputs_url(script_path) if script_path else ""
        result_url = _to_outputs_url(result_path) if result_path else ""

        _save_output_record(
            email=email,
            user_id=user_id,
            notebook_id=notebook_id,
            output_type="podcast",
            file_name=Path(audio_path).name if audio_path else "podcast.wav",
            file_path=audio_path or "",
            result_path=result_path or str(output_dir),
            download_url=audio_url,
        )

        from fastapi_app.kb_records import add_output_record
        try:
            add_output_record(
                user_email=email,
                notebook_id=notebook_id,
                output_type="podcast",
                file_name=Path(audio_path).name if audio_path else "podcast.wav",
                download_url=audio_url
            )
        except Exception as e:
            log.warning("[generate-podcast] failed to write JSON record: %s", e)

        return {
            "success": True,
            "result_path": result_url,
            "audio_path": audio_url,
            "script_path": script_url,
            "output_file_id": f"kb_podcast_{int(time.time())}"
        }

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/generate-mindmap")
async def generate_mindmap_from_kb(
    file_paths: List[str] = Body(..., embed=True),
    user_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    api_url: Optional[str] = Body(None, embed=True),
    api_key: Optional[str] = Body(None, embed=True),
    model: str = Body("deepseek-v3.2", embed=True),
    mindmap_style: str = Body("default", embed=True),
    max_depth: int = Body(3, embed=True),
    language: str = Body("zh", embed=True),
):
    """
    从知识库生成思维导图。支持本地文件与「搜索引入」的 URL：路径用 _resolve_local_path；URL 优先用已存 .md，否则抓取后写临时 .md。
    """
    try:
        api_url, api_key = _require_llm_config(api_url, api_key)
        project_root = get_project_root()
        ts = int(time.time())
        # New layout: outputs/{title}_{id}/mindmap/{ts}/
        if notebook_id:
            paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
            output_dir = paths.feature_output_dir("mindmap", ts)
        else:
            output_dir = _outputs_dir(email, notebook_id, f"{ts}_mindmap_input")
        output_dir.mkdir(parents=True, exist_ok=True)
        local_file_paths: List[str] = []

        for f in (file_paths or []):
            ps = (f or "").strip()
            if ps.startswith("http://") or ps.startswith("https://"):
                content = None
                local_md = _resolve_link_to_local_md(email, notebook_id, ps)
                if local_md is not None:
                    try:
                        content = local_md.read_text(encoding="utf-8", errors="replace")
                        log.info("[generate-mindmap] 网页来源使用已存 .md: %s", local_md.name)
                    except Exception as e:
                        log.warning("[generate-mindmap] 读取已存 .md 失败: %s", e)
                if not (content or "").strip():
                    try:
                        content = fetch_page_text(ps, max_chars=100000)
                    except Exception as e:
                        log.warning("[generate-mindmap] 抓取 URL 失败 %s: %s", ps[:60], e)
                        content = ""
                if (content or "").strip():
                    tmp_md = output_dir / f"link_{len(local_file_paths)}.md"
                    tmp_md.write_text(content.strip(), encoding="utf-8")
                    local_file_paths.append(str(tmp_md))
            else:
                local_path = _resolve_local_path(ps)
                if not local_path.exists() or not local_path.is_file():
                    raise HTTPException(status_code=404, detail=f"File not found: {ps}")
                local_file_paths.append(str(local_path))

        if not local_file_paths:
            raise HTTPException(status_code=400, detail="No valid files provided")

        # Get vector store base directory
        vector_store_base_dir = _vector_store_base_dir(email, notebook_id)

        # Prepare request
        mindmap_req = KBMindMapRequest(
            file_ids=local_file_paths,
            vector_store_base_dir=vector_store_base_dir,
            chat_api_url=api_url,
            api_key=api_key,
            model=model,
            mindmap_style=mindmap_style,
            max_depth=max_depth,
            language=language
        )
        mindmap_req.email = email

        state = KBMindMapState(request=mindmap_req)

        # Run workflow via registry (统一使用 run_workflow)
        result_state = await run_workflow("kb_mindmap", state)

        # Extract results
        mermaid_code = ""
        result_path = ""

        if isinstance(result_state, dict):
            mermaid_code = result_state.get("mermaid_code", "")
            result_path = result_state.get("result_path", "")
        else:
            mermaid_code = getattr(result_state, "mermaid_code", "")
            result_path = getattr(result_state, "result_path", "")

        mindmap_path = ""
        if result_path:
            mmd_path = Path(result_path) / "mindmap.mmd"
            if (not mmd_path.exists()) and mermaid_code:
                try:
                    mmd_path.write_text(mermaid_code, encoding="utf-8")
                except Exception:
                    pass
            if mmd_path.exists():
                mindmap_path = _to_outputs_url(str(mmd_path))

        _save_output_record(
            email=email,
            user_id=user_id,
            notebook_id=notebook_id,
            output_type="mindmap",
            file_name="mindmap.mmd",
            file_path=str(Path(result_path) / "mindmap.mmd") if result_path else "",
            result_path=result_path or "",
            download_url=mindmap_path,
        )

        from fastapi_app.kb_records import add_output_record
        try:
            add_output_record(
                user_email=email,
                notebook_id=notebook_id,
                output_type="mindmap",
                file_name="mindmap.mmd",
                download_url=mindmap_path
            )
        except Exception as e:
            log.warning("[generate-mindmap] failed to write JSON record: %s", e)

        return {
            "success": True,
            "result_path": _to_outputs_url(result_path) if result_path else "",
            "mermaid_code": mermaid_code,
            "mindmap_path": mindmap_path,
            "output_file_id": f"kb_mindmap_{int(time.time())}"
        }

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


def _collect_figure_images(
    mgr: "SourceManager",
    file_paths: List[str],
    project_root: Path,
) -> List[tuple]:
    """
    从 file_paths 对应的 source 中收集 MinerU 提取的 figure 图片。
    返回 [(source_stem, image_path), ...]
    """
    IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"}
    results: List[tuple] = []

    for fp in (file_paths or []):
        ps = (fp or "").strip()
        if ps.startswith("http://") or ps.startswith("https://"):
            continue

        local_path = _resolve_local_path(ps)
        stem = local_path.stem

        # 1) 检查 MinerU auto/images/ 目录下的图片
        mineru_root = mgr.get_mineru_root(stem)
        if mineru_root and mineru_root.exists():
            images_dir = mineru_root / "images"
            scan_dir = images_dir if images_dir.is_dir() else mineru_root
            for img in sorted(scan_dir.iterdir()):
                if img.is_file() and img.suffix.lower() in IMAGE_EXTS:
                    results.append((stem, img))

        # 2) 如果 MinerU 没有图片，检查原始文件本身是否是图片
        if not any(s == stem for s, _ in results):
            if local_path.exists() and local_path.suffix.lower() in IMAGE_EXTS:
                results.append((stem, local_path))

    return results


@router.post("/generate-drawio")
async def generate_drawio_from_kb(
    file_paths: List[str] = Body(..., embed=True),
    user_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    api_url: Optional[str] = Body(None, embed=True),
    api_key: Optional[str] = Body(None, embed=True),
    model: str = Body("deepseek-v3.2", embed=True),
    diagram_type: str = Body("auto", embed=True),
    diagram_style: str = Body("default", embed=True),
    language: str = Body("zh", embed=True),
    text_content: Optional[str] = Body(None, embed=True),
):
    """
    从知识库选中文件生成 DrawIO 图表。

    注意：此功能正在重构中，暂时不可用。
    优先：思维导图生成（/generate-mindmap）、播客生成（/generate-podcast）
    """
    raise HTTPException(
        status_code=501,
        detail="DrawIO 生成功能正在重构中，暂时不可用。请使用思维导图生成功能（/api/v1/kb/generate-mindmap）作为替代。"
    )




@router.post("/save-mindmap")
async def save_mindmap_to_file(
    file_url: str = Body(..., embed=True),
    content: str = Body(..., embed=True),
):
    """
    Save edited Mermaid mindmap code back to the output file.
    """
    try:
        if not file_url:
            raise HTTPException(status_code=400, detail="File URL is required")

        local_path = Path(_from_outputs_url(file_url))
        if not local_path.is_absolute():
            local_path = (get_project_root() / local_path).resolve()

        outputs_root = (get_project_root() / "outputs").resolve()
        try:
            local_path.relative_to(outputs_root)
        except ValueError:
            raise HTTPException(status_code=400, detail="Invalid output path")

        if local_path.suffix.lower() not in {".mmd", ".mermaid", ".md"}:
            raise HTTPException(status_code=400, detail="Invalid mindmap file type")

        local_path.parent.mkdir(parents=True, exist_ok=True)
        local_path.write_text(content or "", encoding="utf-8")

        return {
            "success": True,
            "mindmap_path": _to_outputs_url(str(local_path))
        }
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


# ===================== Flashcard 闪卡 =====================

@router.post("/generate-flashcards")
async def generate_flashcards(
    file_paths: List[str] = Body(..., embed=True),
    email: str = Body(..., embed=True),
    user_id: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    api_url: Optional[str] = Body(None, embed=True),
    api_key: Optional[str] = Body(None, embed=True),
    model: str = Body("deepseek-v3.2", embed=True),
    language: str = Body("zh", embed=True),
    card_count: int = Body(20, embed=True),
):
    """从知识库文件生成闪卡"""
    try:
        api_url, api_key = _require_llm_config(api_url, api_key)
        from fastapi_app.services.flashcard_service import generate_flashcards_with_llm

        local_paths = []
        for f in file_paths:
            ps = (f or "").strip()
            if ps.startswith("http://") or ps.startswith("https://"):
                local_md = _resolve_link_to_local_md(email, notebook_id, ps)
                if local_md and local_md.exists():
                    local_paths.append(str(local_md))
            else:
                local_path = _resolve_local_path(f)
                if local_path.exists():
                    local_paths.append(str(local_path))

        if not local_paths:
            raise HTTPException(status_code=400, detail="No valid files provided")

        text_content = _extract_text_from_files(local_paths, max_chars=50000)
        if not text_content.strip():
            raise HTTPException(status_code=400, detail="No text content extracted")

        log.info("[generate-flashcards] text_len=%d, files=%d", len(text_content), len(local_paths))

        flashcards = await generate_flashcards_with_llm(
            text_content=text_content,
            api_url=api_url,
            api_key=api_key,
            model=model,
            language=language,
            card_count=card_count,
        )
        if not flashcards:
            raise HTTPException(status_code=500, detail="Failed to generate flashcards")

        ts = int(time.time())
        flashcard_set_id = f"flashcard_{ts}"
        if notebook_id:
            paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
            output_dir = paths.feature_output_dir("flashcard", ts)
        else:
            output_dir = _outputs_dir(email, notebook_id, flashcard_set_id)
        output_dir.mkdir(parents=True, exist_ok=True)

        flashcard_data = {
            "id": flashcard_set_id,
            "notebook_id": notebook_id,
            "flashcards": [fc.dict() for fc in flashcards],
            "created_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
            "source_files": file_paths,
            "total_count": len(flashcards),
        }
        (output_dir / "flashcards.json").write_text(
            json.dumps(flashcard_data, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        log.info("[generate-flashcards] 成功生成 %d 张闪卡", len(flashcards))

        from fastapi_app.kb_records import add_output_record
        try:
            add_output_record(
                user_email=email,
                notebook_id=notebook_id,
                output_type="flashcard",
                file_name="flashcards.json",
                download_url=_to_outputs_url(str(output_dir))
            )
        except Exception as e:
            log.warning("[generate-flashcards] failed to write JSON record: %s", e)

        return {
            "success": True,
            "flashcards": [fc.dict() for fc in flashcards],
            "flashcard_set_id": flashcard_set_id,
            "total_count": len(flashcards),
            "result_path": _to_outputs_url(str(output_dir)),
        }
    except HTTPException:
        raise
    except Exception as e:
        log.exception("[generate-flashcards] failed")
        raise HTTPException(status_code=500, detail=str(e))


# ===================== Quiz 测验 =====================

@router.post("/generate-quiz")
async def generate_quiz(
    file_paths: List[str] = Body(..., embed=True),
    email: str = Body(..., embed=True),
    user_id: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    api_url: Optional[str] = Body(None, embed=True),
    api_key: Optional[str] = Body(None, embed=True),
    model: str = Body("deepseek-v3.2", embed=True),
    language: str = Body("en", embed=True),
    question_count: int = Body(10, embed=True),
):
    """生成 Quiz 测验题目"""
    try:
        api_url, api_key = _require_llm_config(api_url, api_key)
        from fastapi_app.services.quiz_service import generate_quiz_with_llm

        local_paths = []
        for f in file_paths:
            ps = (f or "").strip()
            if ps.startswith("http://") or ps.startswith("https://"):
                local_md = _resolve_link_to_local_md(email, notebook_id, ps)
                if local_md and local_md.exists():
                    local_paths.append(str(local_md))
            else:
                local_path = _resolve_local_path(f)
                if local_path.exists():
                    local_paths.append(str(local_path))

        if not local_paths:
            raise HTTPException(status_code=400, detail="No valid files provided")

        text_content = _extract_text_from_files(local_paths, max_chars=50000)
        if not text_content.strip():
            raise HTTPException(status_code=400, detail="No text content extracted")

        log.info("[generate-quiz] text_len=%d, files=%d", len(text_content), len(local_paths))

        questions = await generate_quiz_with_llm(
            text_content=text_content,
            api_url=api_url,
            api_key=api_key,
            model=model,
            language=language,
            question_count=question_count,
        )
        if not questions:
            raise HTTPException(status_code=500, detail="Failed to generate quiz")

        ts = int(time.time())
        quiz_id = f"quiz_{ts}"
        if notebook_id:
            paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
            output_dir = paths.feature_output_dir("quiz", ts)
        else:
            output_dir = _outputs_dir(email, notebook_id, quiz_id)
        output_dir.mkdir(parents=True, exist_ok=True)

        quiz_data = {
            "id": quiz_id,
            "notebook_id": notebook_id,
            "questions": [q.dict() for q in questions],
            "created_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
            "source_files": file_paths,
            "total_count": len(questions),
        }
        (output_dir / "quiz.json").write_text(
            json.dumps(quiz_data, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        log.info("[generate-quiz] 成功生成 %d 道题目", len(questions))

        from fastapi_app.kb_records import add_output_record
        try:
            add_output_record(
                user_email=email,
                notebook_id=notebook_id,
                output_type="quiz",
                file_name="quiz.json",
                download_url=_to_outputs_url(str(output_dir))
            )
        except Exception as e:
            log.warning("[generate-quiz] failed to write JSON record: %s", e)

        return {
            "success": True,
            "questions": [q.dict() for q in questions],
            "quiz_id": quiz_id,
            "total_count": len(questions),
            "result_path": _to_outputs_url(str(output_dir)),
        }
    except HTTPException:
        raise
    except Exception as e:
        log.exception("[generate-quiz] failed")
        raise HTTPException(status_code=500, detail=str(e))


# ===================== Flashcard / Quiz 读取端点 =====================

@router.get("/list-flashcard-sets")
async def list_flashcard_sets(
    notebook_id: str,
    notebook_title: Optional[str] = None,
    user_id: Optional[str] = None,
    email: Optional[str] = None,
):
    """列出某 notebook 下所有已保存的闪卡集合（按时间倒序）"""
    try:
        paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
        flashcard_root = paths.root / "flashcard"
        sets = []
        if flashcard_root.exists():
            for ts_dir in flashcard_root.iterdir():
                if not ts_dir.is_dir():
                    continue
                json_file = ts_dir / "flashcards.json"
                if not json_file.exists():
                    continue
                try:
                    data = json.loads(json_file.read_text(encoding="utf-8"))
                    sets.append({
                        "set_id": ts_dir.name,
                        "id": data.get("id", ""),
                        "created_at": data.get("created_at", ""),
                        "total_count": data.get("total_count", 0),
                        "source_files": data.get("source_files", []),
                    })
                except Exception:
                    continue
        sets.sort(key=lambda x: x["set_id"], reverse=True)
        return {"success": True, "sets": sets}
    except Exception as e:
        log.exception("[list-flashcard-sets] failed")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/list-quiz-sets")
async def list_quiz_sets(
    notebook_id: str,
    notebook_title: Optional[str] = None,
    user_id: Optional[str] = None,
    email: Optional[str] = None,
):
    """列出某 notebook 下所有已保存的测验集合（按时间倒序）"""
    try:
        paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
        quiz_root = paths.root / "quiz"
        sets = []
        if quiz_root.exists():
            for ts_dir in quiz_root.iterdir():
                if not ts_dir.is_dir():
                    continue
                json_file = ts_dir / "quiz.json"
                if not json_file.exists():
                    continue
                try:
                    data = json.loads(json_file.read_text(encoding="utf-8"))
                    sets.append({
                        "set_id": ts_dir.name,
                        "id": data.get("id", ""),
                        "created_at": data.get("created_at", ""),
                        "total_count": data.get("total_count", 0),
                        "source_files": data.get("source_files", []),
                    })
                except Exception:
                    continue
        sets.sort(key=lambda x: x["set_id"], reverse=True)
        return {"success": True, "sets": sets}
    except Exception as e:
        log.exception("[list-quiz-sets] failed")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/get-flashcard-set")
async def get_flashcard_set(
    notebook_id: str,
    set_id: str,
    notebook_title: Optional[str] = None,
    user_id: Optional[str] = None,
    email: Optional[str] = None,
):
    """读取指定闪卡集合的完整数据"""
    try:
        paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
        json_file = paths.root / "flashcard" / set_id / "flashcards.json"
        if not json_file.exists():
            raise HTTPException(status_code=404, detail="Flashcard set not found")
        data = json.loads(json_file.read_text(encoding="utf-8"))
        return {"success": True, **data}
    except HTTPException:
        raise
    except Exception as e:
        log.exception("[get-flashcard-set] failed")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/get-quiz-set")
async def get_quiz_set(
    notebook_id: str,
    set_id: str,
    notebook_title: Optional[str] = None,
    user_id: Optional[str] = None,
    email: Optional[str] = None,
):
    """读取指定测验集合的完整数据"""
    try:
        paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
        json_file = paths.root / "quiz" / set_id / "quiz.json"
        if not json_file.exists():
            raise HTTPException(status_code=404, detail="Quiz set not found")
        data = json.loads(json_file.read_text(encoding="utf-8"))
        return {"success": True, **data}
    except HTTPException:
        raise
    except Exception as e:
        log.exception("[get-quiz-set] failed")
        raise HTTPException(status_code=500, detail=str(e))


# ============================================================================
# DeepResearch Integration
# ============================================================================

@router.post("/deep-research")
async def run_deep_research(
    query: str = Body(..., embed=True),
    notebook_id: str = Body(..., embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    user_id: Optional[str] = Body(None, embed=True),
    email: Optional[str] = Body(None, embed=True),
    max_iterations: int = Body(50, embed=True),
):
    """
    运行 DeepResearch 深度研究并将结果保存为 source

    Args:
        query: 研究问题
        notebook_id: Notebook ID
        notebook_title: Notebook 标题
        user_id: 用户 ID
        email: 用户邮箱
        max_iterations: 最大迭代次数

    Returns:
        {
            "success": bool,
            "query": str,
            "answer": str,
            "source_info": {...},  # 保存的 source 信息
            "error": str (optional)
        }
    """
    try:
        from fastapi_app.services.deep_research_integration import DeepResearchIntegration

        log.info(f"[deep-research] 开始深度研究: {query}")

        # 1. 运行完整的 DeepResearch
        integration = DeepResearchIntegration()
        result = await integration.run_research(
            query=query,
            max_iterations=max_iterations
        )

        if not result["success"]:
            return result

        # 2. 将结果保存为 source
        paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
        mgr = SourceManager(paths)

        # 格式化为 Markdown
        markdown_content = integration.format_result_as_markdown(result)

        # 保存为文本 source
        source_info = await mgr.import_text(
            text=markdown_content,
            title=f"DeepResearch: {query[:50]}"
        )

        log.info(f"[deep-research] 已保存结果: {source_info.original_path}")

        # 3. 自动 embedding
        try:
            vector_base = str(paths.vector_store_dir)
            file_list = [{"path": str(source_info.original_path)}]
            await process_knowledge_base_files(
                file_list=file_list,
                vector_base=vector_base,
                email=email or "default",
                user_id=user_id or "default",
                notebook_id=notebook_id,
            )
            log.info(f"[deep-research] 已完成 embedding")
        except Exception as e:
            log.warning(f"[deep-research] Embedding 失败: {e}")

        return {
            "success": True,
            "query": query,
            "answer": result["answer"],
            "source_info": {
                "file_type": source_info.file_type,
                "original_path": str(source_info.original_path),
                "markdown_path": str(source_info.markdown_path) if source_info.markdown_path else None,
            },
            "sources_count": len(result.get("sources", [])),
        }

    except Exception as e:
        log.exception("[deep-research] 执行失败")
        raise HTTPException(status_code=500, detail=str(e))


# ============================================================================
# Search & Add Integration
# ============================================================================

@router.post("/search-and-add")
async def search_and_add(
    query: str = Body(..., embed=True),
    notebook_id: str = Body(..., embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    user_id: Optional[str] = Body(None, embed=True),
    email: Optional[str] = Body(None, embed=True),
    top_k: int = Body(10, embed=True),
    search_provider: Optional[str] = Body(None, embed=True),
    search_api_key: Optional[str] = Body(None, embed=True),
):
    """
    搜索并爬取 Top K 结果，保存为 source

    Args:
        query: 搜索查询
        notebook_id: Notebook ID
        notebook_title: Notebook 标题
        user_id: 用户 ID
        email: 用户邮箱
        top_k: 返回前 K 个结果
        search_provider: 搜索引擎提供商
        search_api_key: 搜索 API 密钥

    Returns:
        {
            "success": bool,
            "query": str,
            "sources_count": int,
            "crawled_count": int,
            "source_info": {...}
        }
    """
    try:
        from fastapi_app.services.search_and_add_service import SearchAndAddService
        from fastapi_app.kb_records import add_source_record

        log.info(f"[search-and-add] 开始搜索: {query}, top_k={top_k}")

        # 1. 搜索并爬取
        service = SearchAndAddService()
        result = await service.search_and_crawl(
            query=query,
            top_k=top_k,
            search_provider=_resolve_search_provider(search_provider),
            search_api_key=_resolve_search_api_key(search_provider, search_api_key) or None,
        )

        if not result["success"]:
            return result

        sources = result["sources"]
        if not sources:
            return {
                "success": False,
                "query": query,
                "error": "未找到搜索结果"
            }

        # 2. 将所有结果合并为一个 Markdown 文档
        paths = get_notebook_paths(notebook_id, notebook_title or "", email or user_id)
        mgr = SourceManager(paths)

        markdown_content = service.format_sources_as_markdown(sources)

        # 保存为文本 source
        source_info = await mgr.import_text(
            text=markdown_content,
            title=f"Search: {query[:50]}"
        )
        project_root = get_project_root()
        rel = source_info.original_path.relative_to(project_root)
        static_path = "/" + rel.as_posix()

        try:
            add_source_record(
                user_email=email or user_id or DEFAULT_EMAIL,
                notebook_id=notebook_id,
                file_name=source_info.original_path.name,
                file_path=str(source_info.original_path),
                static_url=static_path,
                file_size=source_info.original_path.stat().st_size,
                file_type="text/markdown",
            )
        except Exception as record_err:
            log.warning("[search-and-add] failed to write source record: %s", record_err)

        log.info(f"[search-and-add] 已保存 {len(sources)} 个结果: {source_info.original_path}")

        # 3. 自动 embedding
        try:
            vector_base = str(paths.vector_store_dir)
            file_list = [{"path": str(source_info.original_path)}]
            await process_knowledge_base_files(
                file_list=file_list,
                vector_base=vector_base,
                email=email or "default",
                user_id=user_id or "default",
                notebook_id=notebook_id,
            )
            log.info(f"[search-and-add] 已完成 embedding")
        except Exception as e:
            log.warning(f"[search-and-add] Embedding 失败: {e}")

        crawled_count = sum(1 for s in sources if s["crawl_success"])

        return {
            "success": True,
            "query": query,
            "sources_count": len(sources),
            "crawled_count": crawled_count,
            "source_info": {
                "file_type": source_info.file_type,
                "original_path": str(source_info.original_path),
                "markdown_path": str(source_info.markdown_path) if source_info.markdown_path else None,
            }
        }

    except Exception as e:
        log.exception("[search-and-add] 执行失败")
        raise HTTPException(status_code=500, detail=str(e))
